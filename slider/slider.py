from autohive_integrations_sdk import (
    Integration, ExecutionContext, ActionHandler
)
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
import uuid
import os
import base64
from io import BytesIO
from PIL import Image

# Create the integration using the config.json
slider = Integration.load()

# Global dictionaries to store presentations and uploaded images in memory
presentations = {}
uploaded_images = {}

# Slide layout mapping
LAYOUT_MAP = {
    "title": 0,
    "title_content": 1,
    "section_header": 2,
    "two_content": 3,
    "comparison": 4,
    "title_only": 5,
    "blank": 6,
    "content_with_caption": 7,
    "picture_with_caption": 8
}

# Chart type mapping
CHART_TYPE_MAP = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "area": XL_CHART_TYPE.AREA,
    "xy_scatter": XL_CHART_TYPE.XY_SCATTER
}

def process_files(files: List[Dict[str, Any]]) -> Dict[str, BytesIO]:
    """Process files from the files parameter and return streams by filename"""
    processed_files = {}
    if files:
        for file_item in files:
            # Decode base64 content similar to Gmail implementation
            content_as_string = file_item['content']
            
            # Add necessary padding
            padded_content_string = content_as_string + '=' * (-len(content_as_string) % 4)
            
            # Decode the file content
            file_binary_data = base64.urlsafe_b64decode(padded_content_string.encode('ascii'))
            file_stream = BytesIO(file_binary_data)
            
            processed_files[file_item['name']] = file_stream
    
    return processed_files

def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# ---- Action Handlers ----

@slider.action("create_presentation")
class CreatePresentationAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        title = inputs.get("title")
        subtitle = inputs.get("subtitle")
        files = inputs.get("files", [])
        
        # Process files to get template if provided
        processed_files = process_files(files)
        
        # Create presentation from template or blank
        template_file = None
        for filename, file_stream in processed_files.items():
            if filename.lower().endswith('.pptx'):
                template_file = file_stream
                break
        
        if template_file:
            prs = Presentation(template_file)
        else:
            prs = Presentation()
        
        # Add title slide if title is provided
        if title:
            if len(prs.slides) == 0:
                # Add title slide
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
            else:
                slide = prs.slides[0]
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = title
            
            # Set subtitle if provided
            if subtitle and len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle
        
        # Generate unique ID and store presentation
        presentation_id = str(uuid.uuid4())
        presentations[presentation_id] = prs
        
        return {
            "presentation_id": presentation_id,
            "slide_count": len(prs.slides)
        }

@slider.action("add_slide")
class AddSlideAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        # Always use blank layout (index 6)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        return {
            "slide_index": len(prs.slides) - 1,
            "slide_count": len(prs.slides)
        }

@slider.action("add_text")
class AddTextAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        text = inputs["text"]
        position = inputs["position"]
        formatting = inputs.get("formatting", {})
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        
        # Create text box with size validation
        left = Inches(position["left"])
        top = Inches(position["top"])
        width = Inches(position["width"])
        height = Inches(position["height"])
        
        # Validate text box doesn't exceed slide boundaries
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Adjust if text box exceeds slide boundaries
        if left + width > slide_width:
            width = slide_width - left - Inches(0.1)  # Leave small margin
        if top + height > slide_height:
            height = slide_height - top - Inches(0.1)  # Leave small margin
            
        # Ensure minimum text box size
        if width < Inches(0.5):
            width = Inches(0.5)
        if height < Inches(0.3):
            height = Inches(0.3)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = text
        
        # Apply text containment features to prevent overflow
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Force text to fit shape
        tf.word_wrap = True  # Enable word wrapping
        
        # Set default margins to prevent text touching edges
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        # Apply formatting
        p = tf.paragraphs[0]
        if formatting.get("font_size"):
            p.font.size = Pt(formatting["font_size"])
        if formatting.get("bold"):
            p.font.bold = True
        if formatting.get("italic"):
            p.font.italic = True
        if formatting.get("color"):
            rgb = hex_to_rgb(formatting["color"])
            p.font.color.rgb = RGBColor(*rgb)
        
        return {
            "shape_id": str(txBox.shape_id)
        }

@slider.action("add_image")
class AddImageAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        position = inputs["position"]
        files = inputs.get("files", [])
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        # Process files to get the first image file
        processed_files = process_files(files)
        image_file = None
        
        # Check both filename extensions and content types
        for file_item in files:
            filename = file_item['name']
            content_type = file_item.get('contentType', '')
            
            # Check if it's an image by extension or content type
            is_image_by_extension = any(filename.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'])
            is_image_by_content_type = content_type.startswith('image/')
            
            if is_image_by_extension or is_image_by_content_type:
                image_file = processed_files[filename]
                break
        
        if not image_file:
            raise ValueError("No image file found in files parameter")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        
        # Add image with size validation
        left = Inches(position["left"])
        top = Inches(position["top"])
        
        # Get slide dimensions for boundary checking
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        if "width" in position and "height" in position:
            width = Inches(position["width"])
            height = Inches(position["height"])
            
            # Validate image doesn't exceed slide boundaries
            if left + width > slide_width:
                width = slide_width - left - Inches(0.1)  # Leave small margin
            if top + height > slide_height:
                height = slide_height - top - Inches(0.1)  # Leave small margin
                
            # Ensure minimum image size
            if width < Inches(0.5):
                width = Inches(0.5)
            if height < Inches(0.5):
                height = Inches(0.5)
                
            pic = slide.shapes.add_picture(image_file, left, top, width, height)
        else:
            # If no explicit size, add at original size but validate position
            if left > slide_width - Inches(1):
                left = slide_width - Inches(1)
            if top > slide_height - Inches(1):
                top = slide_height - Inches(1)
                
            pic = slide.shapes.add_picture(image_file, left, top)
        
        return {
            "shape_id": str(pic.shape_id)
        }

@slider.action("add_table")
class AddTableAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        rows = inputs["rows"]
        cols = inputs["cols"]
        position = inputs["position"]
        data = inputs.get("data", [])
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        
        # Create table with size validation
        left = Inches(position["left"])
        top = Inches(position["top"])
        width = Inches(position["width"])
        height = Inches(position["height"])
        
        # Validate table doesn't exceed slide boundaries
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Adjust if table exceeds slide boundaries
        if left + width > slide_width:
            width = slide_width - left - Inches(0.1)  # Leave small margin
        if top + height > slide_height:
            height = slide_height - top - Inches(0.1)  # Leave small margin
            
        # Ensure minimum table size
        if width < Inches(1):
            width = Inches(1)
        if height < Inches(0.5):
            height = Inches(0.5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Fill table with data if provided
        for row_idx, row_data in enumerate(data[:rows]):
            for col_idx, cell_value in enumerate(row_data[:cols]):
                table.cell(row_idx, col_idx).text = str(cell_value)
        
        return {
            "table_id": f"table_{slide_index}_{len(slide.shapes)}"
        }

@slider.action("add_chart")
class AddChartAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        chart_type = inputs["chart_type"]
        position = inputs["position"]
        data = inputs["data"]
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = data["categories"]
        
        for series in data["series"]:
            chart_data.add_series(series["name"], series["values"])
        
        # Add chart with size validation
        left = Inches(position["left"])
        top = Inches(position["top"])
        width = Inches(position["width"])
        height = Inches(position["height"])
        
        # Validate chart doesn't exceed slide boundaries
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Adjust if chart exceeds slide boundaries
        if left + width > slide_width:
            width = slide_width - left - Inches(0.1)  # Leave small margin
        if top + height > slide_height:
            height = slide_height - top - Inches(0.1)  # Leave small margin
            
        # Ensure minimum chart size
        if width < Inches(1):
            width = Inches(1)
        if height < Inches(1):
            height = Inches(1)
        
        chart_type_enum = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart_shape = slide.shapes.add_chart(chart_type_enum, left, top, width, height, chart_data)
        
        return {
            "chart_id": str(chart_shape.shape_id)
        }

@slider.action("extract_text")
class ExtractTextAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        file_path = inputs["file_path"]
        
        if not os.path.exists(file_path):
            raise ValueError(f"File {file_path} not found")
        
        prs = Presentation(file_path)
        slides_text = []
        all_text_parts = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                    
                for paragraph in shape.text_frame.paragraphs:
                    text_run = ''.join(run.text for run in paragraph.runs)
                    if text_run.strip():
                        slide_text.append(text_run.strip())
                        all_text_parts.append(text_run.strip())
            
            slides_text.append({
                "slide_index": slide_idx,
                "text_content": slide_text
            })
        
        return {
            "slides": slides_text,
            "all_text": "\n".join(all_text_parts)
        }

@slider.action("modify_slide")
class ModifySlideAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        updates = inputs["updates"]
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        modified = False
        
        # Update title if provided
        if "title" in updates and slide.shapes.title:
            slide.shapes.title.text = updates["title"]
            modified = True
        
        # Perform text replacements
        if "replace_text" in updates:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                    
                for paragraph in shape.text_frame.paragraphs:
                    for replacement in updates["replace_text"]:
                        old_text = replacement["old_text"]
                        new_text = replacement["new_text"]
                        
                        paragraph_text = ''.join(run.text for run in paragraph.runs)
                        if old_text in paragraph_text:
                            # Replace text in the paragraph
                            new_paragraph_text = paragraph_text.replace(old_text, new_text)
                            
                            # Clear existing runs and add new text
                            paragraph.clear()
                            paragraph.text = new_paragraph_text
                            modified = True
        
        return {
            "modified": modified
        }

@slider.action("save_presentation")
class SavePresentationAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        file_path = inputs["file_path"]
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        
        # Save presentation to memory buffer instead of disk
        try:
            from io import BytesIO
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            file_content = buffer.getvalue()
            
            # Encode as base64
            content_base64 = base64.b64encode(file_content).decode('utf-8')
            
            # Get file name from path
            file_name = os.path.basename(file_path)
            
            return {
                "saved": True,
                "file_path": file_path,
                "file": {
                    "content": content_base64,
                    "name": file_name,
                    "contentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                }
            }
        except Exception as e:
            return {
                "saved": False,
                "file_path": file_path,
                "file": {
                    "content": "",
                    "name": os.path.basename(file_path),
                    "contentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                },
                "error": f"Could not generate presentation for streaming: {str(e)}"
            }

@slider.action("set_text_autosize")
class SetTextAutosizeAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        shape_index = inputs["shape_index"]
        autosize_type = inputs["autosize_type"]
        word_wrap = inputs.get("word_wrap", None)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            raise ValueError(f"Shape index {shape_index} out of range")
        
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            raise ValueError("Shape does not have a text frame")
        
        text_frame = shape.text_frame
        
        # Map string values to MSO_AUTO_SIZE constants
        autosize_map = {
            "NONE": MSO_AUTO_SIZE.NONE,
            "SHAPE_TO_FIT_TEXT": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
            "TEXT_TO_FIT_SHAPE": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        }
        
        if autosize_type not in autosize_map:
            raise ValueError(f"Invalid autosize_type. Must be one of: {list(autosize_map.keys())}")
        
        # Set autosize
        text_frame.auto_size = autosize_map[autosize_type]
        
        # Set word wrap if provided
        if word_wrap is not None:
            text_frame.word_wrap = word_wrap
            
        return {
            "success": True,
            "autosize_type": autosize_type,
            "word_wrap": text_frame.word_wrap
        }

@slider.action("fit_text_to_shape")
class FitTextToShapeAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        shape_index = inputs["shape_index"]
        max_size = inputs.get("max_size", 48)  # Maximum font size in points
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            raise ValueError(f"Shape index {shape_index} out of range")
        
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            raise ValueError("Shape does not have a text frame")
        
        text_frame = shape.text_frame
        
        # Use the fit_text method to automatically size text to fit shape
        text_frame.fit_text(max_size=max_size)
        
        return {
            "success": True,
            "max_size": max_size,
            "auto_size": "TEXT_TO_FIT_SHAPE"
        }

@slider.action("set_text_margins")
class SetTextMarginsAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        shape_index = inputs["shape_index"]
        margins = inputs["margins"]
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            raise ValueError(f"Shape index {shape_index} out of range")
        
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            raise ValueError("Shape does not have a text frame")
        
        text_frame = shape.text_frame
        
        # Set margins (values in inches)
        if "left" in margins:
            text_frame.margin_left = Inches(margins["left"])
        if "right" in margins:
            text_frame.margin_right = Inches(margins["right"])
        if "top" in margins:
            text_frame.margin_top = Inches(margins["top"])
        if "bottom" in margins:
            text_frame.margin_bottom = Inches(margins["bottom"])
            
        return {
            "success": True,
            "margins_set": margins
        }

@slider.action("set_text_alignment")
class SetTextAlignmentAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        shape_index = inputs["shape_index"]
        vertical_anchor = inputs.get("vertical_anchor", None)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            raise ValueError(f"Shape index {shape_index} out of range")
        
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            raise ValueError("Shape does not have a text frame")
        
        text_frame = shape.text_frame
        
        # Map string values to MSO_ANCHOR constants
        anchor_map = {
            "TOP": MSO_ANCHOR.TOP,
            "MIDDLE": MSO_ANCHOR.MIDDLE,
            "BOTTOM": MSO_ANCHOR.BOTTOM
        }
        
        if vertical_anchor and vertical_anchor in anchor_map:
            text_frame.vertical_anchor = anchor_map[vertical_anchor]
            
        return {
            "success": True,
            "vertical_anchor": vertical_anchor
        }

@slider.action("set_slide_background_color")
class SetSlideBackgroundColorAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        color = inputs["color"]
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        background = slide.background
        
        # Set solid color background (this will automatically override master inheritance)
        fill = background.fill
        fill.solid()
        
        # Handle different color input formats
        if isinstance(color, dict):
            if "rgb" in color:
                # RGB format: {"rgb": [255, 0, 0]}
                r, g, b = color["rgb"]
                fill.fore_color.rgb = RGBColor(r, g, b)
            elif "theme_color" in color:
                # Theme color format: {"theme_color": "ACCENT_1"}
                theme_map = {
                    "ACCENT_1": MSO_THEME_COLOR.ACCENT_1,
                    "ACCENT_2": MSO_THEME_COLOR.ACCENT_2,
                    "ACCENT_3": MSO_THEME_COLOR.ACCENT_3,
                    "ACCENT_4": MSO_THEME_COLOR.ACCENT_4,
                    "ACCENT_5": MSO_THEME_COLOR.ACCENT_5,
                    "ACCENT_6": MSO_THEME_COLOR.ACCENT_6,
                    "BACKGROUND_1": MSO_THEME_COLOR.BACKGROUND_1,
                    "BACKGROUND_2": MSO_THEME_COLOR.BACKGROUND_2,
                    "DARK_1": MSO_THEME_COLOR.DARK_1,
                    "DARK_2": MSO_THEME_COLOR.DARK_2,
                    "LIGHT_1": MSO_THEME_COLOR.LIGHT_1,
                    "LIGHT_2": MSO_THEME_COLOR.LIGHT_2
                }
                theme_color = color["theme_color"]
                if theme_color in theme_map:
                    fill.fore_color.theme_color = theme_map[theme_color]
                else:
                    raise ValueError(f"Invalid theme color: {theme_color}")
        elif isinstance(color, str):
            # Hex format: "#FF0000"
            if color.startswith("#"):
                hex_color = color[1:]
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                fill.fore_color.rgb = RGBColor(r, g, b)
            else:
                raise ValueError("String color must be in hex format (e.g., '#FF0000')")
        else:
            raise ValueError("Color must be dict with 'rgb'/'theme_color' or hex string")
            
        return {
            "success": True,
            "color_set": color
        }

@slider.action("set_slide_background_gradient")
class SetSlideBackgroundGradientAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        angle = inputs.get("angle", 90)  # Default 90 degrees (bottom to top)
        gradient_stops = inputs.get("gradient_stops", [])
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        background = slide.background
        
        # Set gradient background (this will automatically override master inheritance)
        fill = background.fill
        fill.gradient()
        
        # Set gradient angle
        fill.gradient_angle = angle
        
        # Set gradient stops if provided
        if gradient_stops:
            stops = fill.gradient_stops
            # Clear existing stops except first two (minimum required)
            while len(stops) > 2:
                stops[-1].remove()
            
            for i, stop_info in enumerate(gradient_stops[:len(stops)]):
                stop = stops[i]
                if "position" in stop_info:
                    stop.position = stop_info["position"]
                if "color" in stop_info:
                    color_info = stop_info["color"]
                    if isinstance(color_info, dict) and "rgb" in color_info:
                        r, g, b = color_info["rgb"]
                        stop.color.rgb = RGBColor(r, g, b)
                    elif isinstance(color_info, dict) and "theme_color" in color_info:
                        theme_map = {
                            "ACCENT_1": MSO_THEME_COLOR.ACCENT_1,
                            "ACCENT_2": MSO_THEME_COLOR.ACCENT_2,
                            "ACCENT_3": MSO_THEME_COLOR.ACCENT_3,
                            "ACCENT_4": MSO_THEME_COLOR.ACCENT_4,
                            "ACCENT_5": MSO_THEME_COLOR.ACCENT_5,
                            "ACCENT_6": MSO_THEME_COLOR.ACCENT_6
                        }
                        theme_color = color_info["theme_color"]
                        if theme_color in theme_map:
                            stop.color.theme_color = theme_map[theme_color]
            
        return {
            "success": True,
            "gradient_angle": angle,
            "gradient_stops_applied": len(gradient_stops) if gradient_stops else 2
        }

@slider.action("add_background_image_workaround")
class AddBackgroundImageWorkaroundAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        files = inputs.get("files", [])
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        # Process files to get the first image file
        processed_files = process_files(files)
        image_file = None
        
        # Check both filename extensions and content types
        for file_item in files:
            filename = file_item['name']
            content_type = file_item.get('contentType', '')
            
            # Check if it's an image by extension or content type
            is_image_by_extension = any(filename.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'])
            is_image_by_content_type = content_type.startswith('image/')
            
            if is_image_by_extension or is_image_by_content_type:
                image_file = processed_files[filename]
                break
        
        if not image_file:
            raise ValueError("No image file found in files parameter")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Add picture that covers entire slide (position at 0,0)
        picture = slide.shapes.add_picture(image_file, 0, 0, slide_width, slide_height)
        
        # Send picture to back so other elements appear on top
        # Note: python-pptx doesn't have a direct "send to back" method
        # The picture will be behind elements added after it
        
        return {
            "success": True,
            "method": "workaround_full_slide_picture",
            "picture_width": slide_width,
            "picture_height": slide_height,
            "note": "Image added as full-slide picture. Add other elements after this for proper layering."
        }

@slider.action("reset_slide_background")
class ResetSlideBackgroundAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        
        # Reset to inherit background from master/layout
        slide.follow_master_background = True
        
        return {
            "success": True,
            "follow_master_background": True,
            "note": "Slide background reset to inherit from master/layout"
        }

@slider.action("add_bullet_list")
class AddBulletListAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        bullet_items = inputs["bullet_items"]
        position = inputs["position"]
        formatting = inputs.get("formatting", {})
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        
        # Create text box with size validation for bullet list
        left = Inches(position["left"])
        top = Inches(position["top"])
        width = Inches(position["width"])
        height = Inches(position["height"])
        
        # Validate bullet list doesn't exceed slide boundaries
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Adjust if bullet list exceeds slide boundaries
        if left + width > slide_width:
            width = slide_width - left - Inches(0.1)  # Leave small margin
        if top + height > slide_height:
            height = slide_height - top - Inches(0.1)  # Leave small margin
            
        # Ensure minimum bullet list size
        if width < Inches(1):
            width = Inches(1)
        if height < Inches(0.5):
            height = Inches(0.5)
        
        # Always create a text box for consistent behavior and avoid double bullets
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        bullet_placeholder = False  # We're not using a placeholder
        
        # Apply text containment features
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        # Add bullet items
        for i, item in enumerate(bullet_items):
            text = item.get("text", "")
            level = item.get("level", 0)
            
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            if bullet_placeholder:
                # Use built-in bullet formatting from placeholder
                p.text = text
                p.level = min(max(level, 0), 8)
            else:
                # Create manual bullets with Unicode symbols only when no placeholder
                bullet_symbols = ["•", "◦", "▪", "▫", "‣"]  # Different levels
                bullet_symbol = bullet_symbols[min(level, len(bullet_symbols)-1)]
                indent = "    " * level  # Manual indentation
                p.text = f"{indent}{bullet_symbol} {text}"
            
            p.alignment = PP_ALIGN.LEFT
        
        # Apply formatting to all paragraphs
        if formatting:
            for paragraph in text_frame.paragraphs:
                if formatting.get("font_size"):
                    paragraph.font.size = Pt(formatting["font_size"])
                if formatting.get("bold"):
                    paragraph.font.bold = formatting["bold"]
                if formatting.get("italic"):
                    paragraph.font.italic = formatting["italic"]
                if formatting.get("color"):
                    rgb = hex_to_rgb(formatting["color"])
                    paragraph.font.color.rgb = RGBColor(*rgb)
        
        shape_id = str(textbox.shape_id)
        
        return {
            "shape_id": shape_id,
            "items_count": len(bullet_items)
        }


