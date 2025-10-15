from autohive_integrations_sdk import (
    Integration, ExecutionContext, ActionHandler
)
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
import uuid
import os
import base64
from io import BytesIO
from PIL import Image
import markdown
from bs4 import BeautifulSoup
import re

# Debug mode for font sizing (set via environment variable)
FONT_SIZE_DEBUG = os.environ.get('FONT_SIZE_DEBUG', 'false').lower() == 'true'

# Load integration from config.json in the same directory as this file
_current_dir = os.path.dirname(os.path.abspath(__file__))
_config_path = os.path.join(_current_dir, "config.json")
slide_maker = Integration.load(_config_path)

presentations = {}
uploaded_images = {}

# Only blank slides are supported (layout index 6)
BLANK_LAYOUT_INDEX = 6

CHART_TYPE_MAP = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "area": XL_CHART_TYPE.AREA,
    "xy_scatter": XL_CHART_TYPE.XY_SCATTER
}

def process_files(files: List[Dict[str, Any]]) -> Dict[str, BytesIO]:
    """Process files from the files parameter and return streams by filename"""
    processed_files = {}
    if files:
        for file_item in files:
            content_as_string = file_item['content']
            
            padded_content_string = content_as_string + '=' * (-len(content_as_string) % 4)
            
            file_binary_data = base64.urlsafe_b64decode(padded_content_string.encode('ascii'))
            file_stream = BytesIO(file_binary_data)
            
            processed_files[file_item['name']] = file_stream
    
    return processed_files

def load_presentation_from_files(presentation_id: str, files: List[Dict[str, Any]]) -> None:
    """Load presentation from files parameter if not in memory"""
    if presentation_id not in presentations and files:
        processed_files = process_files(files)
        
        for filename, file_stream in processed_files.items():
            if filename.lower().endswith('.pptx') or filename.lower().endswith('.bin'):
                try:
                    prs = Presentation(file_stream)
                    presentations[presentation_id] = prs
                    return
                except Exception as e:
                    continue
        
        # If no valid PowerPoint file found, provide better error message
        available_files = list(processed_files.keys())
        raise ValueError(f"No valid PowerPoint file found in files. Tried to load: {available_files}. Files may be corrupted or not PowerPoint format.")
    elif presentation_id not in presentations:
        raise ValueError(f"Presentation {presentation_id} not found and no files provided for loading")

async def save_and_return_presentation(original_result: Dict[str, Any], presentation_id: str, context: ExecutionContext, custom_filename: str = None) -> Dict[str, Any]:
    """Helper to save presentation and return combined result"""
    if presentation_id not in presentations:
        raise ValueError(f"Presentation {presentation_id} not found")

    prs = presentations[presentation_id]

    if custom_filename:
        # Remove any existing .pptx extensions first, then add one
        file_path = custom_filename
        while file_path.lower().endswith('.pptx'):
            file_path = file_path[:-5]  # Remove .pptx
        file_path += '.pptx'  # Add exactly one .pptx
    else:
        file_path = f"{presentation_id}.pptx"

    # Save presentation to memory buffer (inlined from SavePresentationAction)
    try:
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        file_content = buffer.getvalue()

        # Encode as base64
        content_base64 = base64.b64encode(file_content).decode('utf-8')

        # Get file name from path
        file_name = os.path.basename(file_path)

        save_result = {
            "saved": True,
            "file_path": file_path,
            "file": {
                "content": content_base64,
                "name": file_name,
                "contentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            }
        }
    except Exception as e:
        save_result = {
            "saved": False,
            "file_path": file_path,
            "file": {
                "content": "",
                "name": os.path.basename(file_path),
                "contentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            },
            "error": f"Could not generate presentation for streaming: {str(e)}"
        }

    combined_result = original_result.copy()
    combined_result.update({
        "saved": save_result["saved"],
        "file_path": save_result["file_path"],
        "file": save_result["file"],
        "error": save_result.get("error", "")
    })
    return combined_result

def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def calculate_overlap(element1_pos, element2_pos):
    """Calculate overlap between two rectangular elements"""
    left1, top1, right1, bottom1 = element1_pos["left"], element1_pos["top"], element1_pos["right"], element1_pos["bottom"]
    left2, top2, right2, bottom2 = element2_pos["left"], element2_pos["top"], element2_pos["right"], element2_pos["bottom"]
    
    overlap_left = max(left1, left2)
    overlap_top = max(top1, top2)
    overlap_right = min(right1, right2)
    overlap_bottom = min(bottom1, bottom2)
    
    # Check if there's an overlap
    if overlap_left < overlap_right and overlap_top < overlap_bottom:
        overlap_width = overlap_right - overlap_left
        overlap_height = overlap_bottom - overlap_top
        overlap_area = overlap_width * overlap_height
        
        # Calculate percentages of overlap relative to each element
        area1 = element1_pos["width"] * element1_pos["height"]
        area2 = element2_pos["width"] * element2_pos["height"]
        
        overlap_percent1 = (overlap_area / area1 * 100) if area1 > 0 else 0
        overlap_percent2 = (overlap_area / area2 * 100) if area2 > 0 else 0
        
        return {
            "overlaps": True,
            "overlap_area": round(overlap_area, 3),
            "overlap_width": round(overlap_width, 3),
            "overlap_height": round(overlap_height, 3),
            "overlap_position": {
                "left": round(overlap_left, 3),
                "top": round(overlap_top, 3),
                "right": round(overlap_right, 3),
                "bottom": round(overlap_bottom, 3)
            },
            "overlap_percent_element1": round(overlap_percent1, 1),
            "overlap_percent_element2": round(overlap_percent2, 1)
        }
    else:
        return {"overlaps": False}

def analyze_all_element_overlaps(elements):
    """Analyze overlaps between ALL elements on the slide"""
    overlaps = []
    total_overlaps = 0
    
    for i in range(len(elements)):
        for j in range(i + 1, len(elements)):
            element1 = elements[i]
            element2 = elements[j]
            
            overlap_info = calculate_overlap(element1["position"], element2["position"])
            
            if overlap_info["overlaps"]:
                total_overlaps += 1
                overlap_details = {
                    "element1_index": element1["index"],
                    "element1_type": element1["type"],
                    "element1_name": element1.get("name", ""),
                    "element2_index": element2["index"],
                    "element2_type": element2["type"], 
                    "element2_name": element2.get("name", ""),
                    "overlap_area": overlap_info["overlap_area"],
                    "overlap_dimensions": {
                        "width": overlap_info["overlap_width"],
                        "height": overlap_info["overlap_height"]
                    },
                    "overlap_position": overlap_info["overlap_position"],
                    "overlap_percentages": {
                        "element1_coverage": overlap_info["overlap_percent_element1"],
                        "element2_coverage": overlap_info["overlap_percent_element2"]
                    },
                    "severity": "high" if max(overlap_info["overlap_percent_element1"], overlap_info["overlap_percent_element2"]) > 50 else 
                              "medium" if max(overlap_info["overlap_percent_element1"], overlap_info["overlap_percent_element2"]) > 20 else "low",
                    "description": f"{element1['type']} (#{element1['index']}) overlaps {element2['type']} (#{element2['index']}) by {overlap_info['overlap_area']:.2f} sq inches ({max(overlap_info['overlap_percent_element1'], overlap_info['overlap_percent_element2']):.1f}% coverage)"
                }
                overlaps.append(overlap_details)
    
    return overlaps, total_overlaps

def detect_placeholders_with_metadata(text):
    """
    Detect placeholder patterns in text and extract embedded metadata.

    Patterns supported:
    - [Placeholder] - Simple bracket placeholder
    - {Placeholder} - Curly brace placeholder
    - {{Placeholder}} - Double curly placeholder

    Metadata format: [Field, Key=Value, Key=Value]
    Example: [Title, Fontsize=32pt, Bold=true, Font=Sofia Pro, Color=#FF0000]

    Supported metadata keys:
    - Fontsize=32pt - Force specific font size
    - Font=Sofia Pro - Font face to use
    - Bold=true/false - Apply bold formatting
    - Italic=true/false - Apply italic formatting
    - Underline=true/false - Apply underline
    - Color=#HEX or rgb(r,g,b) - Text color
    - Format=... - Hint for data formatting (dates, numbers, etc.)
    - Required=true/false - Validation hint
    - Fallback=text - Default value hint

    Returns:
        tuple: (list of placeholder strings, dict of metadata by placeholder)
    """
    import re

    if not text:
        return [], {}

    patterns = {
        'double_curly': r'\{\{([^\}]+)\}\}',  # {{Placeholder}} - check first (greedy)
        'bracket': r'\[([^\]]+)\]',            # [Placeholder]
        'curly': r'\{([^\}]+)\}'               # {Placeholder}
    }

    placeholders = []
    metadata = {}

    for pattern_type, regex in patterns.items():
        matches = re.findall(regex, text)
        for match in matches:
            # Reconstruct full placeholder
            if pattern_type == 'bracket':
                full_placeholder = f"[{match}]"
            elif pattern_type == 'curly':
                full_placeholder = f"{{{match}}}"
            else:  # double_curly
                full_placeholder = f"{{{{{match}}}}}"

            # Skip if already found (double_curly might overlap with curly)
            if full_placeholder in placeholders:
                continue

            placeholders.append(full_placeholder)

            # Extract metadata if contains commas
            if ',' in match:
                parts = [p.strip() for p in match.split(',')]
                field_name = parts[0]

                meta = {"field": field_name}
                for part in parts[1:]:
                    part_clean = part.strip()

                    if '=' in part_clean:
                        # Explicit key=value format
                        key, value = part_clean.split('=', 1)
                        meta[key.strip().lower()] = value.strip()
                    else:
                        # Check for negation prefix (!)
                        if part_clean.startswith('!'):
                            # Negation: !Bold → Bold=false
                            flag_name = part_clean[1:].strip().lower()
                            if flag_name in ['bold', 'italic', 'underline']:
                                meta[flag_name] = 'false'
                        else:
                            # Shorthand boolean: Bold → Bold=true
                            part_lower = part_clean.lower()
                            if part_lower in ['bold', 'italic', 'underline']:
                                meta[part_lower] = 'true'

                # Only add if has metadata beyond field name
                if len(meta) > 1:
                    metadata[full_placeholder] = meta

    return placeholders, metadata

def strip_conflicting_markdown(text, metadata):
    """
    Strip markdown formatting characters that conflict with metadata formatting.

    If metadata specifies Bold=true, we'll apply bold formatting, so we should
    strip ** and __ from the text to avoid literal asterisks appearing.

    Args:
        text (str): Replacement text that may contain markdown
        metadata (dict): Placeholder metadata (e.g., {'bold': 'true', 'italic': 'false'})

    Returns:
        str: Text with conflicting markdown characters removed
    """
    if not metadata:
        return text

    cleaned = text

    # Strip bold markers if Bold=true in metadata
    if metadata.get('bold', '').lower() == 'true':
        cleaned = cleaned.replace('**', '').replace('__', '')

    # Strip italic markers if Italic=true in metadata
    if metadata.get('italic', '').lower() == 'true':
        # Only strip single * and _, not ** or __ (those are for bold)
        # This is tricky - for simplicity, remove all * and _ if italic specified
        cleaned = cleaned.replace('*', '').replace('_', '')

    # Strip strikethrough if metadata has underline/strikethrough
    # (Less common, but included for completeness)
    cleaned = cleaned.replace('~~', '')

    # Strip code backticks (always, since we're not in a code context)
    cleaned = cleaned.replace('`', '')

    return cleaned

def get_element_info(shape, index: int, slide_width_inches: float, slide_height_inches: float) -> dict:
    """Extract detailed information about a shape/element including boundary checking"""
    from pptx.util import Inches
    
    # Convert EMU to inches for positions and sizes
    left_inches = shape.left / 914400 if shape.left is not None else 0
    top_inches = shape.top / 914400 if shape.top is not None else 0
    width_inches = shape.width / 914400 if shape.width is not None else 0
    height_inches = shape.height / 914400 if shape.height is not None else 0
    
    right_inches = left_inches + width_inches
    bottom_inches = top_inches + height_inches
    
    warnings = []
    boundary_status = "inside"
    
    if left_inches < 0:
        warnings.append(f"Element extends {abs(left_inches):.3f} inches to the left of slide")
        boundary_status = "outside"
    
    if top_inches < 0:
        warnings.append(f"Element extends {abs(top_inches):.3f} inches above slide")
        boundary_status = "outside"
    
    if right_inches > slide_width_inches:
        overhang = right_inches - slide_width_inches
        warnings.append(f"Element extends {overhang:.3f} inches to the right of slide")
        boundary_status = "outside"
    
    if bottom_inches > slide_height_inches:
        overhang = bottom_inches - slide_height_inches
        warnings.append(f"Element extends {overhang:.3f} inches below slide")
        boundary_status = "outside"
    
    # Determine element type
    element_type = "unknown"
    if hasattr(shape, 'shape_type'):
        shape_type = shape.shape_type
        if shape_type == 1:  # MSO_SHAPE_TYPE.AUTO_SHAPE
            element_type = "shape"
        elif shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            element_type = "image"
        elif shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
            element_type = "table"
        elif shape_type == 3:  # MSO_SHAPE_TYPE.CHART
            element_type = "chart"
        elif shape_type == 17:  # MSO_SHAPE_TYPE.TEXT_BOX
            element_type = "text_box"
        elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            element_type = "text_element"
    
    # Extract text content
    content = ""
    has_text = False
    text_overflow_detected = False
    
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        has_text = True
        try:
            text_frame = shape.text_frame
            
            text_parts = []
            for paragraph in text_frame.paragraphs:
                paragraph_text = ''.join(run.text for run in paragraph.runs)
                if paragraph_text.strip():
                    text_parts.append(paragraph_text.strip())
                    
            content = '\n'.join(text_parts)
                        
        except:
            content = ""
    
    # Get shape name if available
    name = ""
    if hasattr(shape, 'name'):
        name = shape.name or ""

    # Detect placeholders and extract metadata
    placeholders = []
    placeholder_metadata = {}
    is_fillable = False

    if has_text and content:
        placeholders, placeholder_metadata = detect_placeholders_with_metadata(content)
        is_fillable = len(placeholders) > 0

        # Debug logging
        if FONT_SIZE_DEBUG and is_fillable:
            print(f"[PLACEHOLDER DETECTION] Element {index}:")
            print(f"  Content: {content[:60]}")
            print(f"  Placeholders: {placeholders}")
            print(f"  Metadata: {placeholder_metadata}")

    result = {
        "index": index,
        "type": element_type,
        "shape_id": str(shape.shape_id) if hasattr(shape, 'shape_id') else str(index),
        "position": {
            "left": round(left_inches, 3),
            "top": round(top_inches, 3),
            "width": round(width_inches, 3),
            "height": round(height_inches, 3),
            "right": round(right_inches, 3),
            "bottom": round(bottom_inches, 3)
        },
        "content": content,
        "has_text": has_text,
        "name": name,
        "boundary_status": boundary_status,
        "boundary_warnings": warnings
    }

    # Add placeholder info if detected
    if is_fillable:
        result["is_fillable"] = True
        result["placeholders"] = placeholders
        if placeholder_metadata:
            result["placeholder_metadata"] = placeholder_metadata

    return result


# ---- Element Type Detection and Markdown Parsing ----

def detect_element_type_from_markdown(markdown_text):
    """
    Auto-detect element type from markdown content.

    Returns: "table", "bullets", or "text"
    """
    if not markdown_text or not isinstance(markdown_text, str):
        return "text"

    # Check for table (| ... | format with separator line)
    if re.search(r'\|.*\|', markdown_text) and re.search(r'\|[-:\s]+\|', markdown_text):
        return "table"

    # Check for lists (-, *, +, or 1., 2., etc.)
    lines = markdown_text.strip().split('\n')
    list_pattern = r'^\s*[-*+]\s+|^\s*\d+\.\s+'
    list_lines = [l for l in lines if re.match(list_pattern, l)]

    # If at least 2 list items, consider it a bullet list
    if len(list_lines) >= 2:
        return "bullets"

    # Default: plain text
    return "text"

def parse_markdown_table(markdown_text):
    """
    Parse markdown table syntax into structured data.

    Input:  | Col1 | Col2 |
            |------|------|
            | A    | B    |

    Output: {"rows": 2, "cols": 2, "data": [["Col1", "Col2"], ["A", "B"]]}
    """
    lines = [l.strip() for l in markdown_text.strip().split('\n') if l.strip()]

    if len(lines) < 2:
        return {"rows": 0, "cols": 0, "data": []}

    table_data = []

    for line in lines:
        # Skip separator lines (|---|---|)
        if re.match(r'\|[-:\s]+\|', line):
            continue

        # Parse row: split by |, strip whitespace
        cells = [cell.strip() for cell in line.split('|')]
        # Remove empty cells from start and end (if line starts/ends with |)
        if cells and cells[0] == '':
            cells = cells[1:]
        if cells and cells[-1] == '':
            cells = cells[:-1]

        if cells:
            table_data.append(cells)

    if not table_data:
        return {"rows": 0, "cols": 0, "data": []}

    rows = len(table_data)
    cols = len(table_data[0]) if table_data else 0

    return {"rows": rows, "cols": cols, "data": table_data}

def parse_markdown_bullets(markdown_text):
    """
    Parse markdown list syntax into bullet items with levels.

    Input:  - Item 1
              - Sub 1
            - Item 2
            1. Numbered item

    Output: [
        {"text": "Item 1", "level": 0},
        {"text": "Sub 1", "level": 1},
        {"text": "Item 2", "level": 0},
        {"text": "Numbered item", "level": 0}
    ]
    """
    lines = markdown_text.split('\n')
    bullet_items = []

    for line in lines:
        # Match list items: -, *, +, or numbered (1., 2., etc.)
        match = re.match(r'^(\s*)([-*+]|\d+\.)\s+(.+)$', line)
        if match:
            indent = match.group(1)
            marker = match.group(2)
            text = match.group(3)

            # Calculate level based on indentation (2 spaces = 1 level)
            level = len(indent) // 2

            bullet_items.append({
                "text": text,
                "level": min(level, 8)  # Cap at level 8
            })

    return bullet_items

def get_current_slide_elements(slide, slide_width_inches, slide_height_inches):
    """
    Get positions of all current elements on slide for overlap detection.

    Returns: [
        {
            "index": 0,
            "position": {
                "left": 1.0, "top": 1.0, "width": 3.0, "height": 2.0,
                "right": 4.0, "bottom": 3.0
            }
        },
        ...
    ]
    """
    elements = []
    for i, shape in enumerate(slide.shapes):
        left = shape.left / 914400 if shape.left is not None else 0
        top = shape.top / 914400 if shape.top is not None else 0
        width = shape.width / 914400 if shape.width is not None else 0
        height = shape.height / 914400 if shape.height is not None else 0

        elements.append({
            "index": i,
            "position": {
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "right": left + width,
                "bottom": top + height
            }
        })
    return elements

def validate_and_adjust_position(requested_pos, existing_elements, slide_dims, auto_position=False):
    """
    Validate position using existing calculate_overlap() function.
    Adjusts for boundaries and optionally finds alternative position if overlap detected.

    Returns: {
      "position": {left, top, width, height},
      "adjusted": boolean,
      "reason": string|null,
      "fits": boolean
    }
    """
    left = requested_pos["left"]
    top = requested_pos["top"]
    width = requested_pos["width"]
    height = requested_pos["height"]
    slide_width = slide_dims["width"]
    slide_height = slide_dims["height"]

    adjusted = False
    reason = None

    # STEP 1: Boundary check (reuse existing pattern)
    if left + width > slide_width:
        width = slide_width - left - 0.1
        adjusted = True
        reason = "Adjusted width to fit slide boundaries"

    if top + height > slide_height:
        height = slide_height - top - 0.1
        adjusted = True
        if reason:
            reason += "; adjusted height to fit slide boundaries"
        else:
            reason = "Adjusted height to fit slide boundaries"

    # Ensure minimum sizes
    if width < 0.5:
        width = 0.5
    if height < 0.3:
        height = 0.3

    # STEP 2: Overlap check using EXISTING calculate_overlap()
    new_pos = {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "right": left + width,
        "bottom": top + height
    }

    has_overlap = False
    for existing in existing_elements:
        overlap = calculate_overlap(new_pos, existing["position"])
        if overlap["overlaps"]:
            has_overlap = True
            break

    if has_overlap:
        if not auto_position:
            return {"fits": False, "reason": "Element overlaps with existing element(s). Enable auto_position to find alternative placement."}

        # STEP 3: Try alternative positions (shift right/down)
        alternatives = [
            (left + 0.5, top),          # Shift right 0.5"
            (left, top + 0.5),          # Shift down 0.5"
            (left + 1.0, top),          # Shift right 1.0"
            (left, top + 1.0),          # Shift down 1.0"
            (left + 0.5, top + 0.5),    # Shift diagonal 0.5"
            (left + 1.5, top),          # Shift right 1.5"
            (left, top + 1.5),          # Shift down 1.5"
            (left + 1.0, top + 1.0),    # Shift diagonal 1.0"
        ]

        for alt_left, alt_top in alternatives:
            # Check if alternative is within bounds
            if alt_left < 0 or alt_top < 0:
                continue
            if alt_left + width > slide_width or alt_top + height > slide_height:
                continue

            alt_pos = {
                "left": alt_left,
                "top": alt_top,
                "width": width,
                "height": height,
                "right": alt_left + width,
                "bottom": alt_top + height
            }

            # Check this alternative for overlaps
            alt_has_overlap = False
            for existing in existing_elements:
                if calculate_overlap(alt_pos, existing["position"])["overlaps"]:
                    alt_has_overlap = True
                    break

            if not alt_has_overlap:
                # Found a valid position!
                return {
                    "position": {"left": alt_left, "top": alt_top, "width": width, "height": height},
                    "adjusted": True,
                    "reason": f"Repositioned to avoid overlap (moved from {left:.1f}, {top:.1f} to {alt_left:.1f}, {alt_top:.1f})",
                    "fits": True
                }

        # No valid alternative found
        return {
            "fits": False,
            "reason": "No valid position found - all attempted positions overlap with existing elements or exceed slide boundaries"
        }

    # No overlap detected
    return {
        "position": {"left": left, "top": top, "width": width, "height": height},
        "adjusted": adjusted,
        "reason": reason,
        "fits": True
    }


# ---- Internal Element Creation Helpers ----

def _create_text_box(slide, text_content, position, prs):
    """
    Core text box creation logic extracted from add_text action.
    Returns the created shape.
    """
    left = Inches(position["left"])
    top = Inches(position["top"])
    width = Inches(position["width"])
    height = Inches(position["height"])

    # Create text box
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()

    # Set margins
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    # Calculate appropriate font size
    width_inches = position["width"]
    height_inches = position["height"]
    has_markdown = has_markdown_formatting(text_content)
    best_fit_size = calculate_best_fit_font_size(text_content, width_inches, height_inches, max_font_size=18, has_formatting=has_markdown, is_bullets=False, font_face='Calibri', debug=FONT_SIZE_DEBUG)

    # Parse and add formatted text
    html = parse_markdown_with_extensions(text_content)
    soup = BeautifulSoup(html, 'html.parser')
    content_element = soup.find('p') or soup
    paragraph = tf.paragraphs[0]
    _add_formatted_text_from_html(paragraph, content_element, base_font='Calibri', base_size=best_fit_size)

    # Enable word wrap
    tf.word_wrap = True

    return txBox

def _create_table(slide, table_data, position, prs):
    """
    Core table creation logic extracted from add_table action.
    table_data format: {"rows": N, "cols": M, "data": [[...], [...]]}
    Returns the created table shape.
    """
    rows = table_data["rows"]
    cols = table_data["cols"]
    data = table_data.get("data", [])

    left = Inches(position["left"])
    top = Inches(position["top"])
    width = Inches(position["width"])
    height = Inches(position["height"])

    # Create table (keep shape reference)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Fill table with data (parse each cell as markdown with auto-sizing)
    for row_idx, row_data in enumerate(data[:rows]):
        for col_idx, cell_value in enumerate(row_data[:cols]):
            cell = table.cell(row_idx, col_idx)
            cell_text = str(cell_value)

            # Get cell dimensions for font size calculation
            cell_width = table.columns[col_idx].width / 914400 if col_idx < len(table.columns) else 2
            cell_height = table.rows[row_idx].height / 914400 if row_idx < len(table.rows) else 0.5

            # Calculate appropriate font size
            has_markdown = has_markdown_formatting(cell_text)
            best_fit_size = calculate_best_fit_font_size(cell_text, cell_width, cell_height, max_font_size=14, has_formatting=has_markdown, is_bullets=False, font_face='Calibri', debug=FONT_SIZE_DEBUG)

            # Parse cell content as markdown
            html = parse_markdown_with_extensions(cell_text)
            soup = BeautifulSoup(html, 'html.parser')
            content_element = soup.find('p') or soup
            cell.text_frame.clear()
            paragraph = cell.text_frame.paragraphs[0]
            _add_formatted_text_from_html(paragraph, content_element, base_font='Calibri', base_size=best_fit_size)

            # Enable word wrap for cells
            cell.text_frame.word_wrap = True

    return table_shape

def _create_bullet_list(slide, bullet_items, position, prs):
    """
    Core bullet list creation logic extracted from add_bullet_list action.
    bullet_items format: [{"text": "...", "level": 0}, ...]
    Returns the created shape.
    """
    left = Inches(position["left"])
    top = Inches(position["top"])
    width = Inches(position["width"])
    height = Inches(position["height"])

    # Create text box for bullets
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()

    # Set margins
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)

    # Calculate appropriate font size based on all bullet text
    width_inches = position["width"]
    height_inches = position["height"]
    all_bullet_text = '\n'.join([item.get("text", "") for item in bullet_items])
    best_fit_size = calculate_best_fit_font_size(all_bullet_text, width_inches, height_inches, max_font_size=18, has_formatting=True, is_bullets=True, font_face='Calibri', debug=FONT_SIZE_DEBUG)

    # Add bullet items with calculated font size
    for i, item in enumerate(bullet_items):
        text = item.get("text", "")
        level = item.get("level", 0)

        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # Create manual bullets with Unicode symbols
        bullet_symbols = ["•", "◦", "▪", "▫", "‣"]
        bullet_symbol = bullet_symbols[min(level, len(bullet_symbols)-1)]
        indent = "    " * level

        # Add bullet symbol
        run = p.add_run()
        run.text = f"{indent}{bullet_symbol} "
        run.font.name = 'Calibri'
        run.font.size = Pt(best_fit_size)

        # Parse text as markdown and add formatted content
        html = parse_markdown_with_extensions(text)
        soup = BeautifulSoup(html, 'html.parser')
        content_element = soup.find('p') or soup
        _add_formatted_text_from_html(p, content_element, base_font='Calibri', base_size=best_fit_size)

    # Enable word wrap
    text_frame.word_wrap = True

    return textbox


# ---- HTML/Markdown Parser Utilities ----

def get_slide_dimensions_inches(prs):
    """Get slide dimensions in inches"""
    return {
        "width": prs.slide_width / 914400,
        "height": prs.slide_height / 914400
    }

def has_markdown_formatting(text: str) -> bool:
    """Check if text contains markdown formatting markers"""
    if not text:
        return False
    markers = ['**', '*', '`', '~~', '__']
    return any(marker in text for marker in markers)

def get_font_path(font_face):
    """
    Resolve font file path from font name.

    Strategy:
    1. Check if font exists in local fonts/ directory
    2. If not found, check if font has Google Fonts alternative (from font_alternatives.json)
    3. If alternative exists, try to download from Google Fonts
    4. If all fails, return None (triggers heuristic fallback)

    Args:
        font_face (str): Font name (e.g., 'Calibri', 'Sofia Pro Light', etc.)

    Returns:
        str: Full path to font file, or None if not found
    """
    # Map font names to local files (case-insensitive)
    font_files = {
        'calibri': 'Calibri.ttf',
        'calibri bold': 'Calibri-Bold.ttf',
        'arial': 'Arial.ttf',
        'times new roman': 'Times.ttf',
        'times': 'Times.ttf',
        'courier new': 'Courier.ttf',
        'courier': 'Courier.ttf',
    }

    font_key = font_face.lower()
    current_dir = os.path.dirname(os.path.abspath(__file__))
    fonts_dir = os.path.join(current_dir, 'fonts')

    # STEP 1: Check local bundled fonts
    font_file = font_files.get(font_key)
    if font_file:
        font_path = os.path.join(fonts_dir, font_file)
        if os.path.exists(font_path):
            return font_path

    # STEP 2: Try Google Fonts alternative
    try:
        # Import here to avoid circular dependencies
        import sys
        fonts_module_path = os.path.join(current_dir, 'fonts')
        if fonts_module_path not in sys.path:
            sys.path.insert(0, fonts_module_path)

        from google_fonts_downloader import load_font_alternatives, download_google_font

        # Load alternatives mapping
        alternatives = load_font_alternatives()
        google_alternative = alternatives.get(font_key)

        if google_alternative:
            # Parse variant from alternative (e.g., "Poppins Light" → family="Poppins", variant="light")
            parts = google_alternative.split()
            if len(parts) > 1 and parts[-1].lower() in ['thin', 'extralight', 'light', 'regular', 'medium', 'semibold', 'bold', 'extrabold', 'black']:
                family = ' '.join(parts[:-1])
                variant = parts[-1].lower()
            else:
                family = google_alternative
                variant = 'regular'

            # Try to download
            downloaded_path = download_google_font(family, variant)
            if downloaded_path:
                return downloaded_path

    except Exception as e:
        # Google Fonts download failed - continue to fallback
        pass

    # STEP 3: Not found - return None (triggers heuristic fallback)
    return None

def calculate_best_fit_font_size_pillow(text, width_inches, height_inches, max_font_size=18, has_formatting=True, is_bullets=False, font_face='Calibri', debug=False):
    """
    Calculate optimal font size using Pillow for accurate text measurement.

    IMPORTANT: This function only scales DOWN, never UP. If text fits at max_font_size,
    it returns max_font_size (the requested size). Only if it doesn't fit does it
    reduce the size until it fits.

    Args:
        text (str): Text to measure
        width_inches (float): Box width in inches
        height_inches (float): Box height in inches
        max_font_size (int): Requested/maximum font size (default: 18)
        has_formatting (bool): Whether text has bold/italic (affects width)
        is_bullets (bool): Whether this is a bullet list (needs extra space)
        font_face (str): Font name (default: 'Calibri')
        debug (bool): Enable debug logging (default: False)

    Returns:
        int: Font size in points (will be <= max_font_size), or None if font file not available
    """
    from PIL import ImageFont
    import math

    # Get font file path - use BOLD font if text has formatting
    if has_formatting and font_face.lower() == 'calibri':
        # Try bold font for more accurate measurement
        current_dir = os.path.dirname(os.path.abspath(__file__))
        bold_font_path = os.path.join(current_dir, 'fonts', 'Calibri-Bold.ttf')
        if os.path.exists(bold_font_path):
            font_path = bold_font_path
            if debug:
                print(f"  [PILLOW] Using BOLD font for measurement (has_formatting=True)")
        else:
            font_path = get_font_path(font_face)
    else:
        font_path = get_font_path(font_face)

    if not font_path:
        if debug:
            print(f"[PILLOW] Font file not found for '{font_face}', falling back to heuristic")
        return None  # Font file not available, caller will use heuristic

    if debug:
        print(f"\n[PILLOW] Starting calculation:")
        print(f"  Text: \"{text[:50]}{'...' if len(text) > 50 else ''}\"")
        print(f"  Text length: {len(text)} chars")
        print(f"  Box: {width_inches:.2f}w x {height_inches:.2f}h inches")
        print(f"  Max font size requested: {max_font_size}pt")
        print(f"  Font: {font_face} {'(BOLD)' if has_formatting and 'Bold' in font_path else ''}")

    # Handle empty or invalid text
    if not text or len(text.strip()) == 0:
        return max_font_size

    # Remove markdown markers for measurement
    clean_text = text.replace('**', '').replace('*', '').replace('__', '').replace('`', '').replace('~~', '')

    # Convert inches to pixels (72 DPI standard)
    box_width_px = width_inches * 72
    box_height_px = height_inches * 72

    # Calculate usable space (with margins + PowerPoint internal padding)
    # Explicit margins: 0.1" LR, 0.05" TB
    # Additional safety: 0.05" each side for PowerPoint internal padding
    usable_width_px = (box_width_px - 0.3 * 72)  # 0.3 inch total (0.2 explicit + 0.1 safety)
    usable_height_px = (box_height_px - 0.2 * 72)  # 0.2 inch total (0.1 explicit + 0.1 safety)

    if debug:
        print(f"  Box dimensions: {box_width_px:.0f}px x {box_height_px:.0f}px")
        print(f"  Usable space (after margins): {usable_width_px:.0f}px x {usable_height_px:.0f}px")

    # Minimum readable size
    min_font_size = 10

    # Try from requested size down to minimum (only scale DOWN, never UP)
    for font_size in range(max_font_size, min_font_size - 1, -1):  # Step by -1pt for precision
        try:
            # Load font at current size
            font = ImageFont.truetype(font_path, size=font_size)

            # ACTUAL multi-line measurement instead of estimation
            from PIL import Image, ImageDraw
            import textwrap

            # Create dummy image for text measurement
            dummy_img = Image.new('RGB', (1, 1))
            draw = ImageDraw.Draw(dummy_img)

            # Calculate average character width for wrapping
            avg_char_bbox = font.getbbox('M', anchor='lt')  # Use 'M' as average width char
            avg_char_width = avg_char_bbox[2] - avg_char_bbox[0]

            # Calculate how many chars fit per line
            chars_per_line = int(usable_width_px / avg_char_width) if avg_char_width > 0 else 50

            # Wrap text respecting explicit newlines (paragraph breaks)
            # Split by \n first, then wrap each segment
            paragraphs = clean_text.split('\n')
            wrapped_lines = []
            for para in paragraphs:
                if para.strip():  # Non-empty paragraph
                    para_wrapped = textwrap.wrap(para, width=max(chars_per_line, 10))
                    if para_wrapped:
                        wrapped_lines.extend(para_wrapped)
                    else:
                        wrapped_lines.append(para)  # Keep even if doesn't wrap
                else:
                    # Empty line (paragraph break)
                    wrapped_lines.append('')  # Preserve blank lines

            if not wrapped_lines:
                wrapped_lines = [clean_text]  # Fallback

            # Measure ACTUAL height of wrapped text using multiline_textbbox
            wrapped_text = '\n'.join(wrapped_lines)
            multiline_bbox = draw.multiline_textbbox(
                (0, 0),
                wrapped_text,
                font=font,
                spacing=0  # No extra spacing - font already includes line height
            )

            actual_text_height_px = multiline_bbox[3] - multiline_bbox[1]
            actual_line_count = len(wrapped_lines)

            # Small safety buffer: 5% or 10px, whichever is larger
            safety_buffer_px = max(actual_text_height_px * 0.05, 10)
            total_height_needed = actual_text_height_px + safety_buffer_px

            # Check if text fits VERTICALLY
            if total_height_needed <= usable_height_px:
                if debug:
                    fill_percentage = (total_height_needed / usable_height_px) * 100
                    print(f"  [PILLOW] Found fit at {font_size}pt:")
                    print(f"    Wrapped lines: {actual_line_count}")
                    print(f"    Measured height: {actual_text_height_px:.0f}px")
                    print(f"    Safety buffer: {safety_buffer_px:.0f}px")
                    print(f"    Total needed: {total_height_needed:.0f}px")
                    print(f"    Available: {usable_height_px:.0f}px")
                    print(f"    Fill rate: {fill_percentage:.1f}%")
                return font_size  # Return first size that fits (requested size or smaller)

        except Exception as e:
            # If PIL fails for this font size, continue trying smaller sizes
            if debug and font_size == max_font_size:
                print(f"  [PILLOW] Error at {font_size}pt: {e}")
            continue

    if debug:
        print(f"  [PILLOW] WARNING: Even at minimum {min_font_size}pt, text doesn't fit!")
        print(f"    This box is too small for this amount of text.")

    return min_font_size  # Return minimum if nothing fits

def calculate_best_fit_font_size(text, width_inches, height_inches, max_font_size=18, has_formatting=True, is_bullets=False, font_face='Calibri', debug=False):
    """
    Calculate optimal font size to fit text in box.

    Uses Pillow for accurate measurement when font files available.
    Falls back to heuristic approach if font file not found.

    Args:
        text (str): Text to measure
        width_inches (float): Box width in inches
        height_inches (float): Box height in inches
        max_font_size (int): Maximum font size (default: 18)
        has_formatting (bool): Whether text has bold/italic
        is_bullets (bool): Whether this is a bullet list
        font_face (str): Font name (default: 'Calibri')
        debug (bool): Enable debug logging (default: False)

    Returns:
        int: Optimal font size in points
    """

    # Try Pillow-based accurate measurement
    font_size = calculate_best_fit_font_size_pillow(
        text, width_inches, height_inches, max_font_size,
        has_formatting, is_bullets, font_face, debug
    )

    if font_size is not None:
        if debug:
            print(f"  [RESULT] Using PILLOW method: {font_size}pt")
        return font_size

    # Fallback to heuristic (font file not available)
    if debug:
        print(f"  [FALLBACK] Using HEURISTIC method (font file not available)")

    result = calculate_best_fit_font_size_heuristic(
        text, width_inches, height_inches, max_font_size,
        has_formatting, is_bullets, debug
    )

    if debug:
        print(f"  [RESULT] Using HEURISTIC method: {result}pt")

    return result


def calculate_best_fit_font_size_heuristic(text, width_inches, height_inches, max_font_size=18, has_formatting=True, is_bullets=False, debug=False):
    """
    Original heuristic-based font size calculation (fallback).
    Used when font files are not available.

    Algorithm:
    1. Calculate how many characters fit per line based on width
    2. Calculate how many lines needed for total text
    3. Calculate how many lines fit vertically based on height
    4. If lines needed > lines that fit, scale down font size

    Args:
        has_formatting: True if text has bold/italic (takes more space)
        is_bullets: True if this is a bullet list (needs extra space for bullets)
        debug: Enable debug logging
    """
    if not text or width_inches <= 0 or height_inches <= 0:
        return max_font_size

    if debug:
        print(f"\n[HEURISTIC] Starting calculation:")
        print(f"  Text: \"{text[:50]}{'...' if len(text) > 50 else ''}\"")
        print(f"  Box: {width_inches:.2f}w x {height_inches:.2f}h inches")
        print(f"  Max font: {max_font_size}pt")

    # Remove markdown markers for length calculation
    clean_text = text.replace('**', '').replace('*', '').replace('__', '').replace('`', '').replace('~~', '')
    char_count = len(clean_text)

    # Estimate characters per line at max font size
    # At 18pt Calibri: ~8-10 chars per inch width
    # Bold/italic text is significantly wider, especially at large sizes
    if has_formatting:
        # Bold text gets exponentially wider at larger sizes
        # At 32pt bold: ~3-4 chars per inch
        # At 18pt bold: ~5 chars per inch
        if max_font_size >= 28:
            chars_per_inch_width = 3.5  # Large bold titles
        elif max_font_size >= 20:
            chars_per_inch_width = 4.5  # Medium bold headings
        else:
            chars_per_inch_width = 5.5  # Small bold text
    else:
        chars_per_inch_width = 8  # Plain text

    # Bullet lists need extra space for bullet symbols and indentation
    if is_bullets:
        chars_per_inch_width = chars_per_inch_width * 0.8  # 20% reduction for bullets

    usable_width = width_inches - 0.2  # Account for margins
    chars_per_line = usable_width * chars_per_inch_width * (max_font_size / 18)

    if chars_per_line <= 0:
        return max_font_size

    # Estimate lines needed
    lines_needed = char_count / chars_per_line

    # Estimate line height at max font size
    # At 18pt: ~0.25 inches per line
    # Bold text and bullets need slightly more vertical space
    line_height_multiplier = 1.15 if (has_formatting or is_bullets) else 1.0
    line_height_inches = 0.25 * (max_font_size / 18) * line_height_multiplier
    usable_height = height_inches - 0.1  # Account for margins
    max_lines_that_fit = usable_height / line_height_inches

    if lines_needed <= max_lines_that_fit:
        # Text fits at max size
        if debug:
            print(f"  [HEURISTIC] Text fits at max size: {max_font_size}pt")
            print(f"    Lines needed: {lines_needed:.1f}, can fit: {max_lines_that_fit:.1f}")
        return max_font_size

    # Calculate scaling factor based on vertical overflow
    scale_factor = max_lines_that_fit / lines_needed
    calculated_size = max_font_size * scale_factor

    # Add safety margin (reduce by 10% to ensure it fits)
    calculated_size = calculated_size * 0.9

    # Minimum readable size (matches Pillow minimum for consistency)
    min_size = 10
    final_size = max(min_size, min(calculated_size, max_font_size))

    if debug:
        print(f"  [HEURISTIC] Calculated size: {int(final_size)}pt")
        print(f"    Lines needed: {lines_needed:.1f}, can fit: {max_lines_that_fit:.1f}")
        print(f"    Scale factor: {scale_factor:.2f}, with safety: {calculated_size:.1f}pt")

    return int(final_size)

def _add_formatted_text_from_html(paragraph, html_element, base_font='Calibri', base_size=None, base_color=None):
    """
    Parse HTML inline formatting and apply to PowerPoint paragraph.
    Adapted from doc_maker.py for PowerPoint runs.
    Supports: <strong>, <b>, <em>, <i>, <u>, <code>, <s>, <del>

    Args:
        base_font: Uniform font name for all runs (required for auto-sizing)
        base_size: Base font size in points (None = let auto-sizing handle it)
        base_color: RGB tuple (r, g, b) for text color (None = default black)
    """
    for content in html_element.contents:
        if hasattr(content, 'name') and content.name:
            # This is an HTML tag
            if content.name in ['strong', 'b']:
                run = paragraph.add_run()
                run.text = content.get_text()
                run.font.name = base_font
                run.font.bold = True
                if base_size:
                    run.font.size = Pt(base_size)
                if base_color:
                    run.font.color.rgb = RGBColor(*base_color)
            elif content.name in ['em', 'i']:
                run = paragraph.add_run()
                run.text = content.get_text()
                run.font.name = base_font
                run.font.italic = True
                if base_size:
                    run.font.size = Pt(base_size)
                if base_color:
                    run.font.color.rgb = RGBColor(*base_color)
            elif content.name == 'u':
                run = paragraph.add_run()
                run.text = content.get_text()
                run.font.name = base_font
                run.font.underline = True
                if base_size:
                    run.font.size = Pt(base_size)
                if base_color:
                    run.font.color.rgb = RGBColor(*base_color)
            elif content.name == 'code':
                run = paragraph.add_run()
                run.text = content.get_text()
                run.font.name = 'Courier New'  # Code uses monospace
                if base_size:
                    run.font.size = Pt(base_size)
                if base_color:
                    run.font.color.rgb = RGBColor(*base_color)
            elif content.name in ['s', 'del', 'strike']:
                # Note: python-pptx doesn't support strikethrough
                run = paragraph.add_run()
                run.text = content.get_text()
                run.font.name = base_font
                if base_size:
                    run.font.size = Pt(base_size)
                if base_color:
                    run.font.color.rgb = RGBColor(*base_color)
            else:
                # Nested elements - recursively process
                if hasattr(content, 'contents'):
                    _add_formatted_text_from_html(paragraph, content, base_font, base_size, base_color)
                else:
                    # Just add the text content
                    text = content.get_text()
                    if text:
                        run = paragraph.add_run()
                        run.text = text
                        run.font.name = base_font
                        if base_size:
                            run.font.size = Pt(base_size)
                        if base_color:
                            run.font.color.rgb = RGBColor(*base_color)
        else:
            # This is plain text (NavigableString)
            text = str(content)
            if text:
                run = paragraph.add_run()
                run.text = text
                run.font.name = base_font
                if base_size:
                    run.font.size = Pt(base_size)
                if base_color:
                    run.font.color.rgb = RGBColor(*base_color)

def parse_markdown_with_extensions(text, include_tables=False):
    """
    Parse markdown with custom preprocessing for underline.

    Supported formatting:
    - **bold** → Bold
    - *italic* → Italic
    - __underline__ → Underline (custom, preprocessed)
    - `code` → Monospace (Courier New)

    Note: ~~strikethrough~~ is NOT supported (python-pptx API limitation)
    """
    # Preprocess: Convert __text__ to <u>text</u> for underline
    # This must happen BEFORE markdown processes it as bold
    text = re.sub(r'__(.*?)__', r'<u>\1</u>', text)

    # Now convert markdown to HTML
    extensions = ['nl2br']
    if include_tables:
        extensions.extend(['tables', 'fenced_code'])

    html = markdown.markdown(text, extensions=extensions)

    return html

def _add_title_from_html(slide, h1_element, slide_dims):
    """Add H1 as slide title (top, centered, bold)"""
    left = 0.5
    top = 0.5
    width = slide_dims["width"] - 1.0
    height = 1.5

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.clear()

    # Set margins first
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # Calculate appropriate font size (titles are always bold)
    title_text = h1_element.get_text()
    best_fit_size = calculate_best_fit_font_size(title_text, width, height, max_font_size=32, has_formatting=True, is_bullets=False, font_face='Calibri', debug=FONT_SIZE_DEBUG)

    # Parse inline HTML formatting with calculated font size
    paragraph = tf.paragraphs[0]
    _add_formatted_text_from_html(paragraph, h1_element, base_font='Calibri', base_size=best_fit_size)

    # Apply title styling
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.CENTER

    # Enable word wrap
    tf.word_wrap = True

    return txBox, height + 0.5  # Return shape and space used

def _add_heading_from_html(slide, heading_element, current_top, margin_left, content_width, level=2):
    """Add H2/H3/H4 as section heading"""
    # Max font sizes based on heading level
    max_font_sizes = {2: 24, 3: 20, 4: 18, 5: 16, 6: 14}
    max_font_size = max_font_sizes.get(level, 18)

    height = 0.6

    txBox = slide.shapes.add_textbox(
        Inches(margin_left), Inches(current_top),
        Inches(content_width), Inches(height)
    )
    tf = txBox.text_frame
    tf.clear()

    # Set margins
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # Calculate appropriate font size (headings are always bold)
    heading_text = heading_element.get_text()
    best_fit_size = calculate_best_fit_font_size(heading_text, content_width, height, max_font_size=max_font_size, has_formatting=True, is_bullets=False, font_face='Calibri', debug=FONT_SIZE_DEBUG)

    # Parse inline HTML formatting with calculated font size
    paragraph = tf.paragraphs[0]
    _add_formatted_text_from_html(paragraph, heading_element, base_font='Calibri', base_size=best_fit_size)

    # Apply heading styling
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.LEFT

    # Enable word wrap
    tf.word_wrap = True

    return txBox, height + 0.2  # Return shape and space used

def _add_paragraph_from_html(slide, p_element, current_top, margin_left, content_width):
    """Add paragraph as body text"""
    height = 0.8  # Initial estimate

    txBox = slide.shapes.add_textbox(
        Inches(margin_left), Inches(current_top),
        Inches(content_width), Inches(height)
    )
    tf = txBox.text_frame
    tf.clear()

    # Set margins
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # Calculate appropriate font size (check if paragraph has formatting)
    paragraph_text = p_element.get_text()
    has_formatting = len(p_element.find_all(['strong', 'b', 'em', 'i', 'u', 'code'])) > 0
    best_fit_size = calculate_best_fit_font_size(paragraph_text, content_width, height, max_font_size=18, has_formatting=has_formatting, is_bullets=False, font_face='Calibri', debug=FONT_SIZE_DEBUG)

    # Parse inline HTML formatting with calculated font size
    paragraph = tf.paragraphs[0]
    _add_formatted_text_from_html(paragraph, p_element, base_font='Calibri', base_size=best_fit_size)

    # Apply body text styling
    paragraph.alignment = PP_ALIGN.LEFT

    # Enable word wrap
    tf.word_wrap = True

    # Calculate actual height used
    actual_height = txBox.height / 914400

    return txBox, actual_height + 0.2  # Return shape and space used

def _add_bullet_list_from_html(slide, list_element, current_top, margin_left, content_width, numbered=False):
    """Add HTML list as PowerPoint bullet/numbered list"""
    items = []
    for li in list_element.find_all('li', recursive=False):
        items.append(li.get_text().strip())

    if not items:
        return None, 0

    # Estimate height based on number of items
    height = 0.4 * len(items) + 0.5

    txBox = slide.shapes.add_textbox(
        Inches(margin_left + 0.3), Inches(current_top),  # Indent bullets slightly
        Inches(content_width - 0.3), Inches(height)
    )
    tf = txBox.text_frame
    tf.clear()

    # Set margins first
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    # Calculate appropriate font size for all bullets
    all_text = '\n'.join(items)
    best_fit_size = calculate_best_fit_font_size(all_text, content_width - 0.3, height, max_font_size=18, has_formatting=True, is_bullets=True, font_face='Calibri', debug=FONT_SIZE_DEBUG)

    # Add bullet items with calculated font size
    for i, item_text in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Add bullet symbol and text with uniform font
        if numbered:
            run = p.add_run()
            run.text = f"{i+1}. {item_text}"
            run.font.name = 'Calibri'
            run.font.size = Pt(best_fit_size)
        else:
            run = p.add_run()
            run.text = f"• {item_text}"
            run.font.name = 'Calibri'
            run.font.size = Pt(best_fit_size)

        p.alignment = PP_ALIGN.LEFT

    # Enable word wrap
    tf.word_wrap = True

    # Calculate actual height used
    actual_height = txBox.height / 914400

    return txBox, actual_height + 0.3  # Return shape and space used

def _add_table_from_html_pptx(slide, table_element, current_top, margin_left, content_width):
    """
    Add HTML table as PowerPoint table.
    Adapted from doc_maker.py:546-578 for PowerPoint.
    """
    rows = table_element.find_all('tr')
    if not rows:
        return None, 0

    # Determine table dimensions
    max_cols = 0
    for row in rows:
        cols = len(row.find_all(['td', 'th']))
        max_cols = max(max_cols, cols)

    if max_cols == 0:
        return None, 0

    # Estimate table height (0.4 inches per row)
    height = len(rows) * 0.4

    # Create PowerPoint table
    table_shape = slide.shapes.add_table(
        len(rows), max_cols,
        Inches(margin_left), Inches(current_top),
        Inches(content_width), Inches(height)
    )
    table = table_shape.table

    # Fill table data
    for row_idx, row in enumerate(rows):
        cells = row.find_all(['td', 'th'])
        for col_idx, cell in enumerate(cells):
            if col_idx < max_cols:
                ppt_cell = table.cell(row_idx, col_idx)
                text = cell.get_text().strip()
                ppt_cell.text = text

                # Apply formatting
                if ppt_cell.text_frame and ppt_cell.text_frame.paragraphs:
                    for paragraph in ppt_cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)

                        # Make header cells bold
                        if cell.name == 'th':
                            paragraph.font.bold = True

    return table_shape, height + 0.3  # Return shape and space used

def _add_blockquote_from_html(slide, blockquote_element, current_top, margin_left, content_width):
    """Add blockquote as indented text box"""
    text = blockquote_element.get_text().strip()

    height = 0.8
    indent = 0.5

    txBox = slide.shapes.add_textbox(
        Inches(margin_left + indent), Inches(current_top),
        Inches(content_width - indent), Inches(height)
    )
    tf = txBox.text_frame
    tf.clear()

    # Parse inline HTML formatting with uniform font
    paragraph = tf.paragraphs[0]
    _add_formatted_text_from_html(paragraph, blockquote_element, base_font='Calibri', base_size=16)

    # Apply blockquote styling (italic, slightly smaller)
    paragraph.font.italic = True
    paragraph.alignment = PP_ALIGN.LEFT

    # Enable word wrap (no auto_size - we calculated font size already)
    tf.word_wrap = True

    # Calculate actual height used
    actual_height = txBox.height / 914400

    return txBox, actual_height + 0.2  # Return shape and space used

def _add_code_block_from_html(slide, pre_element, current_top, margin_left, content_width):
    """Add code block as monospace text box"""
    code_text = pre_element.get_text()

    # Estimate height based on number of lines
    num_lines = code_text.count('\n') + 1
    height = num_lines * 0.25 + 0.3

    txBox = slide.shapes.add_textbox(
        Inches(margin_left), Inches(current_top),
        Inches(content_width), Inches(height)
    )
    tf = txBox.text_frame
    tf.clear()

    paragraph = tf.paragraphs[0]
    paragraph.text = code_text

    # Apply code block styling (monospace, smaller)
    paragraph.font.name = 'Courier New'
    paragraph.alignment = PP_ALIGN.LEFT

    # Enable word wrap (no auto_size - we calculated font size already)
    tf.word_wrap = True

    # Calculate actual height used
    actual_height = txBox.height / 914400

    return txBox, actual_height + 0.2  # Return shape and space used

def parse_markdown_to_slide(slide, markdown_text, prs):
    """
    Parse markdown → HTML → slide elements with smart layout.

    Similar to doc_maker.py:463-511 but for PowerPoint.
    Returns list of created shapes and elements.
    """
    # Convert markdown to HTML (with strikethrough, underline, tables, and code blocks)
    html = parse_markdown_with_extensions(markdown_text, include_tables=True)
    soup = BeautifulSoup(html, 'html.parser')

    # Get slide dimensions for layout calculations
    slide_dims = get_slide_dimensions_inches(prs)

    # Track vertical position as we add elements
    current_top = 0.5
    margin_left = 0.5
    content_width = slide_dims["width"] - 1.0

    created_shapes = []

    # Process each HTML element in order
    for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'table', 'blockquote', 'pre']):
        shape = None
        space_used = 0
        element_type = element.name

        try:
            if element.name == 'h1':
                # H1 = Main title (top, centered, 32pt)
                shape, space_used = _add_title_from_html(slide, element, slide_dims)
                current_top = 2.2  # Reserve top space for title

            elif element.name == 'h2':
                # H2 = Section heading (24pt, bold)
                shape, space_used = _add_heading_from_html(slide, element, current_top, margin_left, content_width, level=2)
                current_top += space_used

            elif element.name == 'h3':
                # H3 = Sub-heading (20pt, bold)
                shape, space_used = _add_heading_from_html(slide, element, current_top, margin_left, content_width, level=3)
                current_top += space_used

            elif element.name in ['h4', 'h5', 'h6']:
                # H4/H5/H6 = Smaller headings
                level = int(element.name[1])
                shape, space_used = _add_heading_from_html(slide, element, current_top, margin_left, content_width, level=level)
                current_top += space_used

            elif element.name == 'p':
                # Paragraph = Body text (18pt)
                shape, space_used = _add_paragraph_from_html(slide, element, current_top, margin_left, content_width)
                current_top += space_used

            elif element.name == 'ul':
                # Unordered list = Bullet list
                shape, space_used = _add_bullet_list_from_html(slide, element, current_top, margin_left, content_width, numbered=False)
                current_top += space_used

            elif element.name == 'ol':
                # Ordered list = Numbered list
                shape, space_used = _add_bullet_list_from_html(slide, element, current_top, margin_left, content_width, numbered=True)
                current_top += space_used

            elif element.name == 'table':
                # Table = PowerPoint table
                shape, space_used = _add_table_from_html_pptx(slide, element, current_top, margin_left, content_width)
                current_top += space_used

            elif element.name == 'blockquote':
                # Blockquote = Indented text box
                shape, space_used = _add_blockquote_from_html(slide, element, current_top, margin_left, content_width)
                current_top += space_used

            elif element.name == 'pre':
                # Code block = Monospace text box
                shape, space_used = _add_code_block_from_html(slide, element, current_top, margin_left, content_width)
                current_top += space_used

            if shape:
                created_shapes.append({
                    "type": element_type,
                    "shape_id": str(shape.shape_id)
                })

        except Exception as e:
            # Skip elements that fail to process, continue with others
            continue

    return created_shapes


# ---- Action Handlers ----

@slide_maker.action("create_presentation")
class CreatePresentationAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        title = inputs.get("title")
        subtitle = inputs.get("subtitle")
        files = inputs.get("files", [])
        custom_filename = inputs.get("custom_filename")
        
        processed_files = process_files(files)
        
        template_file = None
        for filename, file_stream in processed_files.items():
            if filename.lower().endswith('.pptx'):
                template_file = file_stream
                break
        
        if template_file:
            prs = Presentation(template_file)
        else:
            prs = Presentation()
        
        # Always ensure at least one slide exists
        if len(prs.slides) == 0:
            blank_slide_layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]
            slide = prs.slides.add_slide(blank_slide_layout)
        else:
            slide = prs.slides[0]
        
        # Add title and subtitle as text boxes if provided
        if title:
            
            # Add title as a text box on the blank slide
            title_left = 0.5 
            title_top = 0.5   
            title_width = prs.slide_width / 914400 - 1.0 
            title_height = 1.0 
            
            title_box = slide.shapes.add_textbox(
                Inches(title_left), Inches(title_top),
                Inches(title_width), Inches(title_height)
            )
            title_frame = title_box.text_frame
            title_frame.clear()

            # Calculate appropriate title font size (titles are bold)
            has_markdown = has_markdown_formatting(title)
            best_title_size = calculate_best_fit_font_size(title, title_width, title_height, max_font_size=32, has_formatting=True, is_bullets=False, font_face='Calibri', debug=FONT_SIZE_DEBUG)

            # Always parse title as markdown with calculated font size
            html = parse_markdown_with_extensions(title)
            soup = BeautifulSoup(html, 'html.parser')
            content_element = soup.find('p') or soup
            paragraph = title_frame.paragraphs[0]
            _add_formatted_text_from_html(paragraph, content_element, base_font='Calibri', base_size=best_title_size)

            # Apply standard title styling
            paragraph.font.bold = True

            # Enable word wrap so title can wrap to multiple lines
            title_frame.word_wrap = True

            # Add subtitle if provided
            if subtitle:
                subtitle_top = title_top + title_height + 0.2
                subtitle_height = 0.8

                subtitle_box = slide.shapes.add_textbox(
                    Inches(title_left), Inches(subtitle_top),
                    Inches(title_width), Inches(subtitle_height)
                )
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.clear()

                # Calculate appropriate subtitle font size
                has_markdown = has_markdown_formatting(subtitle)
                best_subtitle_size = calculate_best_fit_font_size(subtitle, title_width, subtitle_height, max_font_size=18, has_formatting=has_markdown, is_bullets=False, font_face='Calibri', debug=FONT_SIZE_DEBUG)

                # Always parse subtitle as markdown with calculated font size
                html = parse_markdown_with_extensions(subtitle)
                soup = BeautifulSoup(html, 'html.parser')
                content_element = soup.find('p') or soup
                paragraph = subtitle_frame.paragraphs[0]
                _add_formatted_text_from_html(paragraph, content_element, base_font='Calibri', base_size=best_subtitle_size)

                # Enable word wrap for subtitle
                subtitle_frame.word_wrap = True
        
        # Generate unique ID and store presentation
        presentation_id = str(uuid.uuid4())
        presentations[presentation_id] = prs
        
        result = {
            "presentation_id": presentation_id,
            "slide_count": len(prs.slides)
        }
        
        return await save_and_return_presentation(result, presentation_id, context, custom_filename)

@slide_maker.action("add_slide")
class AddSlideAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        slide_layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]
        slide = prs.slides.add_slide(slide_layout)
        
        original_result = {
            "slide_index": len(prs.slides) - 1,
            "slide_count": len(prs.slides)
        }
        return await save_and_return_presentation(original_result, presentation_id, context)

@slide_maker.action("add_image")
class AddImageAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        position = inputs["position"]
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        processed_files = process_files(files)
        image_file = None
        
        for file_item in files:
            filename = file_item['name']
            content_type = file_item.get('contentType', '')
            
            # Check if it's an image by extension or content type
            is_image_by_extension = any(filename.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'])
            is_image_by_content_type = content_type.startswith('image/')
            
            if is_image_by_extension or is_image_by_content_type:
                image_file = processed_files[filename]
                break
        
        if not image_file:
            raise ValueError("No image file found in files parameter")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        
        left = Inches(position["left"])
        top = Inches(position["top"])
        
        # Get slide dimensions for boundary checking
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        if "width" in position and "height" in position:
            width = Inches(position["width"])
            height = Inches(position["height"])
            
            # Validate image doesn't exceed slide boundaries
            if left + width > slide_width:
                width = slide_width - left - Inches(0.1)
            if top + height > slide_height:
                height = slide_height - top - Inches(0.1)
                
            if width < Inches(0.5):
                width = Inches(0.5)
            if height < Inches(0.5):
                height = Inches(0.5)
                
            pic = slide.shapes.add_picture(image_file, left, top, width, height)
        else:
            # If no explicit size, add at original size but validate position
            if left > slide_width - Inches(1):
                left = slide_width - Inches(1)
            if top > slide_height - Inches(1):
                top = slide_height - Inches(1)
                
            pic = slide.shapes.add_picture(image_file, left, top)
        
        original_result = {"shape_id": str(pic.shape_id)}
        return await save_and_return_presentation(original_result, presentation_id, context)

@slide_maker.action("add_chart")
class AddChartAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        chart_type = inputs["chart_type"]
        position = inputs["position"]
        data = inputs["data"]
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        
        chart_data = CategoryChartData()
        chart_data.categories = data["categories"]
        
        for series in data["series"]:
            chart_data.add_series(series["name"], series["values"])
        
        left = Inches(position["left"])
        top = Inches(position["top"])
        width = Inches(position["width"])
        height = Inches(position["height"])
        
        # Validate chart doesn't exceed slide boundaries
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Adjust if chart exceeds slide boundaries
        if left + width > slide_width:
            width = slide_width - left - Inches(0.1)  
        if top + height > slide_height:
            height = slide_height - top - Inches(0.1)  
            
        if width < Inches(1):
            width = Inches(1)
        if height < Inches(1):
            height = Inches(1)
        
        chart_type_enum = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart_shape = slide.shapes.add_chart(chart_type_enum, left, top, width, height, chart_data)
        
        original_result = {"chart_id": str(chart_shape.shape_id)}
        return await save_and_return_presentation(original_result, presentation_id, context)

@slide_maker.action("set_text_autosize")
class SetTextAutosizeAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        shape_index = inputs["shape_index"]
        autosize_type = inputs["autosize_type"]
        word_wrap = inputs.get("word_wrap", None)
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            if len(slide.shapes) == 0:
                raise ValueError(f"Shape index {shape_index} out of range. Slide has no elements.")
            else:
                raise ValueError(f"Shape index {shape_index} out of range. Valid range: 0-{len(slide.shapes)-1} ({len(slide.shapes)} elements total).")
        
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            raise ValueError("Shape does not have a text frame")
        
        text_frame = shape.text_frame
        
        # Map string values to MSO_AUTO_SIZE constants
        autosize_map = {
            "NONE": MSO_AUTO_SIZE.NONE,
            "SHAPE_TO_FIT_TEXT": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
            "TEXT_TO_FIT_SHAPE": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        }
        
        if autosize_type not in autosize_map:
            raise ValueError(f"Invalid autosize_type. Must be one of: {list(autosize_map.keys())}")
        
        text_frame.auto_size = autosize_map[autosize_type]
        
        if word_wrap is not None:
            text_frame.word_wrap = word_wrap
            
        # Force recalculation by slightly adjusting text box size
        if autosize_type != "NONE":
            original_width = shape.width
            original_height = shape.height
            
            shape.width = original_width + 1
            shape.height = original_height + 1
            
            shape.width = original_width
            shape.height = original_height
            
        original_result = {
            "success": True,
            "autosize_type": autosize_type,
            "word_wrap": text_frame.word_wrap
        }
        return await save_and_return_presentation(original_result, presentation_id, context)

@slide_maker.action("set_text_margins")
class SetTextMarginsAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        shape_index = inputs["shape_index"]
        margins = inputs["margins"]
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            if len(slide.shapes) == 0:
                raise ValueError(f"Shape index {shape_index} out of range. Slide has no elements.")
            else:
                raise ValueError(f"Shape index {shape_index} out of range. Valid range: 0-{len(slide.shapes)-1} ({len(slide.shapes)} elements total).")
        
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            raise ValueError("Shape does not have a text frame")
        
        text_frame = shape.text_frame
        
        # Set margins (values in inches)
        if "left" in margins:
            text_frame.margin_left = Inches(margins["left"])
        if "right" in margins:
            text_frame.margin_right = Inches(margins["right"])
        if "top" in margins:
            text_frame.margin_top = Inches(margins["top"])
        if "bottom" in margins:
            text_frame.margin_bottom = Inches(margins["bottom"])
            
        result = {
            "success": True,
            "margins_set": margins
        }
        
        return await save_and_return_presentation(result, presentation_id, context)

@slide_maker.action("set_text_alignment")
class SetTextAlignmentAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        shape_index = inputs["shape_index"]
        vertical_anchor = inputs.get("vertical_anchor", None)
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            if len(slide.shapes) == 0:
                raise ValueError(f"Shape index {shape_index} out of range. Slide has no elements.")
            else:
                raise ValueError(f"Shape index {shape_index} out of range. Valid range: 0-{len(slide.shapes)-1} ({len(slide.shapes)} elements total).")
        
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            raise ValueError("Shape does not have a text frame")
        
        text_frame = shape.text_frame
        
        # Map string values to MSO_ANCHOR constants
        anchor_map = {
            "TOP": MSO_ANCHOR.TOP,
            "MIDDLE": MSO_ANCHOR.MIDDLE,
            "BOTTOM": MSO_ANCHOR.BOTTOM
        }
        
        if vertical_anchor and vertical_anchor in anchor_map:
            text_frame.vertical_anchor = anchor_map[vertical_anchor]
            
        result = {
            "success": True,
            "vertical_anchor": vertical_anchor
        }
        
        return await save_and_return_presentation(result, presentation_id, context)

@slide_maker.action("set_slide_background_color")
class SetSlideBackgroundColorAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        color = inputs["color"]
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        background = slide.background
        
        # Set solid color background (this will automatically override master inheritance)
        fill = background.fill
        fill.solid()
        
        # Handle different color input formats
        if isinstance(color, dict):
            if "rgb" in color:
                # RGB format: {"rgb": [255, 0, 0]}
                r, g, b = color["rgb"]
                fill.fore_color.rgb = RGBColor(r, g, b)
            elif "theme_color" in color:
                # Theme color format: {"theme_color": "ACCENT_1"}
                theme_map = {
                    "ACCENT_1": MSO_THEME_COLOR.ACCENT_1,
                    "ACCENT_2": MSO_THEME_COLOR.ACCENT_2,
                    "ACCENT_3": MSO_THEME_COLOR.ACCENT_3,
                    "ACCENT_4": MSO_THEME_COLOR.ACCENT_4,
                    "ACCENT_5": MSO_THEME_COLOR.ACCENT_5,
                    "ACCENT_6": MSO_THEME_COLOR.ACCENT_6,
                    "BACKGROUND_1": MSO_THEME_COLOR.BACKGROUND_1,
                    "BACKGROUND_2": MSO_THEME_COLOR.BACKGROUND_2,
                    "DARK_1": MSO_THEME_COLOR.DARK_1,
                    "DARK_2": MSO_THEME_COLOR.DARK_2,
                    "LIGHT_1": MSO_THEME_COLOR.LIGHT_1,
                    "LIGHT_2": MSO_THEME_COLOR.LIGHT_2
                }
                theme_color = color["theme_color"]
                if theme_color in theme_map:
                    fill.fore_color.theme_color = theme_map[theme_color]
                else:
                    raise ValueError(f"Invalid theme color: {theme_color}")
        elif isinstance(color, str):
            # Hex format: "#FF0000"
            if color.startswith("#"):
                hex_color = color[1:]
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                fill.fore_color.rgb = RGBColor(r, g, b)
            else:
                raise ValueError("String color must be in hex format (e.g., '#FF0000')")
        else:
            raise ValueError("Color must be dict with 'rgb'/'theme_color' or hex string")
            
        result = {
            "success": True,
            "color_set": color
        }
        
        return await save_and_return_presentation(result, presentation_id, context)

@slide_maker.action("set_slide_background_gradient")
class SetSlideBackgroundGradientAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        angle = inputs.get("angle", 90)  # Default 90 degrees (bottom to top)
        gradient_stops = inputs.get("gradient_stops", [])
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        background = slide.background
        
        # Set gradient background (this will automatically override master inheritance)
        fill = background.fill
        fill.gradient()
        
        fill.gradient_angle = angle
        
        # Set gradient stops if provided
        if gradient_stops:
            stops = fill.gradient_stops
            # Clear existing stops except first two (minimum required)
            while len(stops) > 2:
                stops[-1].remove()
            
            for i, stop_info in enumerate(gradient_stops[:len(stops)]):
                stop = stops[i]
                if "position" in stop_info:
                    stop.position = stop_info["position"]
                if "color" in stop_info:
                    color_info = stop_info["color"]
                    if isinstance(color_info, dict) and "rgb" in color_info:
                        r, g, b = color_info["rgb"]
                        stop.color.rgb = RGBColor(r, g, b)
                    elif isinstance(color_info, dict) and "theme_color" in color_info:
                        theme_map = {
                            "ACCENT_1": MSO_THEME_COLOR.ACCENT_1,
                            "ACCENT_2": MSO_THEME_COLOR.ACCENT_2,
                            "ACCENT_3": MSO_THEME_COLOR.ACCENT_3,
                            "ACCENT_4": MSO_THEME_COLOR.ACCENT_4,
                            "ACCENT_5": MSO_THEME_COLOR.ACCENT_5,
                            "ACCENT_6": MSO_THEME_COLOR.ACCENT_6
                        }
                        theme_color = color_info["theme_color"]
                        if theme_color in theme_map:
                            stop.color.theme_color = theme_map[theme_color]
            
        result = {
            "success": True,
            "gradient_angle": angle,
            "gradient_stops_applied": len(gradient_stops) if gradient_stops else 2
        }
        
        return await save_and_return_presentation(result, presentation_id, context)

@slide_maker.action("add_background_image_workaround")
class AddBackgroundImageWorkaroundAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        processed_files = process_files(files)
        image_file = None
        
        # Check both filename extensions and content types
        for file_item in files:
            filename = file_item['name']
            content_type = file_item.get('contentType', '')
            
            # Check if it's an image by extension or content type
            is_image_by_extension = any(filename.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'])
            is_image_by_content_type = content_type.startswith('image/')
            
            if is_image_by_extension or is_image_by_content_type:
                image_file = processed_files[filename]
                break
        
        if not image_file:
            raise ValueError("No image file found in files parameter")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Add picture that covers entire slide (position at 0,0)
        picture = slide.shapes.add_picture(image_file, 0, 0, slide_width, slide_height)
        
        # Send picture to back so other elements appear on top
        # Note: python-pptx doesn't have a direct "send to back" method
        # The picture will be behind elements added after it
        
        result = {
            "success": True,
            "method": "workaround_full_slide_picture",
            "picture_width": slide_width,
            "picture_height": slide_height,
            "note": "Image added as full-slide picture. Add other elements after this for proper layering."
        }
        
        return await save_and_return_presentation(result, presentation_id, context)

@slide_maker.action("reset_slide_background")
class ResetSlideBackgroundAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        
        # Reset background by removing custom background fill
        # This will cause the slide to inherit from master/layout
        background = slide.background
        fill = background.fill
        
        # Clear the background fill to revert to master
        try:
            # Remove any custom background by setting it to no fill
            if hasattr(fill, '_fill') and fill._fill is not None:
                fill._fill.clear()
        except:
            # Alternative approach: set background to inherit from master
            try:
                slide._element.cSld.attrib.pop('showMasterSp', None)
                if hasattr(slide._element.cSld, 'bg'):
                    bg_element = slide._element.cSld.bg
                    if bg_element is not None:
                        slide._element.cSld.remove(bg_element)
            except:
                pass  # If reset fails, continue
        
        result = {
            "success": True,
            "follow_master_background": True,
            "note": "Slide background reset to inherit from master/layout"
        }
        
        return await save_and_return_presentation(result, presentation_id, context)

@slide_maker.action("delete_element")
class DeleteElementAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        shape_index = inputs["shape_index"]
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            if len(slide.shapes) == 0:
                raise ValueError(f"Shape index {shape_index} out of range. Slide has no elements.")
            else:
                raise ValueError(f"Shape index {shape_index} out of range. Valid range: 0-{len(slide.shapes)-1} ({len(slide.shapes)} elements total).")
        
        # Get shape information before deletion for reporting
        shape = slide.shapes[shape_index]
        element_type = "unknown"
        
        # Try to identify the element type
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type
            if shape_type == 1:  # MSO_SHAPE_TYPE.AUTO_SHAPE
                element_type = "shape"
            elif shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                element_type = "image"
            elif shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                element_type = "table"
            elif shape_type == 3:  # MSO_SHAPE_TYPE.CHART
                element_type = "chart"
            elif shape_type == 17:  # MSO_SHAPE_TYPE.TEXT_BOX
                element_type = "text_box"
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                element_type = "text_element"
        
        # Delete the shape
        # In python-pptx, we need to delete by getting the shape's element and removing it from its parent
        shape_element = shape.element
        shape_element.getparent().remove(shape_element)
        
        # Get remaining shape count after deletion
        remaining_shapes = len(slide.shapes)
        
        result = {
            "deleted": True,
            "element_type": element_type,
            "remaining_shapes": remaining_shapes
        }
        
        return await save_and_return_presentation(result, presentation_id, context)

@slide_maker.action("get_slide_elements")
class GetSlideElementsAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs.get("slide_index")  # Now optional
        include_content = inputs.get("include_content", False)
        files = inputs.get("files", [])

        load_presentation_from_files(presentation_id, files)

        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")

        prs = presentations[presentation_id]

        # If no slide_index provided, return data for ALL slides
        if slide_index is None:
            return self._get_all_slides_elements(prs, include_content)

        # Original single-slide logic
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")

        return self._get_single_slide_elements(prs, slide_index, include_content)

    def _get_single_slide_elements(self, prs, slide_index, include_content):
        """Get elements for a single slide"""
        slide = prs.slides[slide_index]

        # Get slide dimensions in inches
        slide_width_inches = prs.slide_width / 914400
        slide_height_inches = prs.slide_height / 914400

        # Get information about all elements on the slide
        elements = []
        elements_outside_boundary = 0
        all_warnings = []

        for i, shape in enumerate(slide.shapes):
            element_info = get_element_info(shape, i, slide_width_inches, slide_height_inches)
            elements.append(element_info)

            # Track boundary violations
            if element_info["boundary_status"] == "outside":
                elements_outside_boundary += 1
                for warning in element_info["boundary_warnings"]:
                    all_warnings.append(f"Element {i} ({element_info['type']}): {warning}")

        # Analyze element overlaps
        overlaps, total_overlaps = analyze_all_element_overlaps(elements)

        # Create summary warning messages
        boundary_warning = ""
        if elements_outside_boundary > 0:
            boundary_warning = f"{elements_outside_boundary} element(s) extend beyond slide boundaries (slide size: {slide_width_inches:.1f}\" x {slide_height_inches:.1f}\")"

        overlap_warning = ""
        if total_overlaps > 0:
            high_severity = len([o for o in overlaps if o["severity"] == "high"])
            if high_severity > 0:
                overlap_warning = f"{total_overlaps} element overlap(s) detected ({high_severity} high severity)"
            else:
                overlap_warning = f"{total_overlaps} element overlap(s) detected"

        combined_warnings = []
        if boundary_warning:
            combined_warnings.append(boundary_warning)
        if overlap_warning:
            combined_warnings.append(overlap_warning)

        overall_warning = "; ".join(combined_warnings) if combined_warnings else ""

        # Create optimized result with minimal tokens
        has_issues = elements_outside_boundary > 0 or total_overlaps > 0

        result = {
            "slide_index": slide_index,
            "total_elements": len(elements),
            "layout_status": "issues_detected" if has_issues else "no_issues",
            "slide_dimensions": {
                "width": round(slide_width_inches, 3),
                "height": round(slide_height_inches, 3)
            }
        }

        # Only include issue details if there are actual issues
        if has_issues:
            result["elements_outside_boundary"] = elements_outside_boundary
            result["total_overlapping_pairs"] = total_overlaps

            if overlaps:
                result["element_overlaps"] = overlaps

        # Optimize element data based on include_content parameter
        optimized_elements = []
        for element in elements:
            optimized_element = {
                "index": element["index"],
                "type": element["type"],
                "shape_id": element["shape_id"],
                "position": {
                    "left": element["position"]["left"],
                    "top": element["position"]["top"],
                    "width": element["position"]["width"],
                    "height": element["position"]["height"]
                }
            }

            # Include content if requested OR if there are issues with this element
            has_element_issues = (element["boundary_status"] != "inside" or
                                element["boundary_warnings"])

            if include_content or has_element_issues:
                if element["has_text"] and element["content"]:
                    optimized_element["content"] = element["content"]

            # Only include problem indicators if there are actual problems
            if element["boundary_status"] != "inside":
                optimized_element["boundary_status"] = element["boundary_status"]

            if element["boundary_warnings"]:
                optimized_element["boundary_warnings"] = element["boundary_warnings"]

            # Include placeholder detection fields if present
            if element.get("is_fillable"):
                optimized_element["is_fillable"] = element["is_fillable"]
                optimized_element["placeholders"] = element["placeholders"]
                if element.get("placeholder_metadata"):
                    optimized_element["placeholder_metadata"] = element["placeholder_metadata"]

            optimized_elements.append(optimized_element)

        result["elements"] = optimized_elements

        return result

    def _get_all_slides_elements(self, prs, include_content):
        """Get elements for all slides in the presentation"""
        slide_width_inches = prs.slide_width / 914400
        slide_height_inches = prs.slide_height / 914400

        all_slides = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_data = self._get_single_slide_elements(prs, slide_idx, include_content)
            all_slides.append(slide_data)

        # Create summary
        total_slides = len(prs.slides)
        slides_with_issues = sum(1 for s in all_slides if s["layout_status"] == "issues_detected")
        total_elements = sum(s["total_elements"] for s in all_slides)

        return {
            "total_slides": total_slides,
            "slides_with_issues": slides_with_issues,
            "total_elements": total_elements,
            "slide_dimensions": {
                "width": round(slide_width_inches, 3),
                "height": round(slide_height_inches, 3)
            },
            "slides": all_slides
        }

@slide_maker.action("reposition_element")
class RepositionElementAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        element_index = inputs["element_index"]
        new_position = inputs.get("position")
        table_cell_updates = inputs.get("table_cell_updates")
        files = inputs.get("files", [])
        
        load_presentation_from_files(presentation_id, files)
        
        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")
        
        slide = prs.slides[slide_index]
        if element_index >= len(slide.shapes):
            if len(slide.shapes) == 0:
                raise ValueError(f"Element index {element_index} out of range. Slide has no elements.")
            else:
                raise ValueError(f"Element index {element_index} out of range. Valid range: 0-{len(slide.shapes)-1} ({len(slide.shapes)} elements total).")
        
        shape = slide.shapes[element_index]
        changes_made = []
        
        # Get element type for reporting
        element_type = "unknown"
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type
            if shape_type == 1:
                element_type = "shape"
            elif shape_type == 13:
                element_type = "image"
            elif shape_type == 19:
                element_type = "table"
            elif shape_type == 3:
                element_type = "chart"
            elif shape_type == 17:
                element_type = "text_box"
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                element_type = "text_element"
        
        # Modify position and size if provided
        if new_position:
            if "left" in new_position:
                shape.left = Inches(new_position["left"])
                changes_made.append(f"Updated left position to {new_position['left']} inches")
            if "top" in new_position:
                shape.top = Inches(new_position["top"])
                changes_made.append(f"Updated top position to {new_position['top']} inches")
            if "width" in new_position:
                shape.width = Inches(new_position["width"])
                changes_made.append(f"Updated width to {new_position['width']} inches")
            if "height" in new_position:
                shape.height = Inches(new_position["height"])
                changes_made.append(f"Updated height to {new_position['height']} inches")

        # Handle table cell updates if this is a table element
        if element_type == "table" and table_cell_updates and hasattr(shape, 'table'):
            table = shape.table
            for cell_update in table_cell_updates:
                row = cell_update.get("row")
                col = cell_update.get("col")
                text = cell_update.get("text")
                cell_formatting = cell_update.get("formatting", {})

                if row is not None and col is not None and row < len(table.rows) and col < len(table.columns):
                    cell = table.cell(row, col)
                    if text is not None:
                        # Clear existing content
                        cell.text_frame.clear()

                        # Get cell dimensions for font size calculation
                        cell_width = table.columns[col].width / 914400 if col < len(table.columns) else 2
                        cell_height = table.rows[row].height / 914400 if row < len(table.rows) else 0.5

                        # Calculate appropriate font size for cell
                        has_markdown = has_markdown_formatting(text)
                        best_fit_size = calculate_best_fit_font_size(text, cell_width, cell_height, max_font_size=14, has_formatting=has_markdown, is_bullets=False, font_face='Calibri', debug=FONT_SIZE_DEBUG)

                        # Always parse as markdown with calculated font size
                        html = parse_markdown_with_extensions(text)
                        soup = BeautifulSoup(html, 'html.parser')

                        # Find the main content element
                        content_element = soup.find('p') or soup

                        # Apply formatted text to cell with calculated font size
                        paragraph = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else cell.text_frame.add_paragraph()
                        _add_formatted_text_from_html(paragraph, content_element, base_font='Calibri', base_size=best_fit_size)

                        # Enable word wrap for cells
                        cell.text_frame.word_wrap = True

                        changes_made.append(f"Updated table cell ({row},{col})")
        
        # Get updated position information
        new_left = shape.left / 914400 if shape.left is not None else 0
        new_top = shape.top / 914400 if shape.top is not None else 0
        new_width = shape.width / 914400 if shape.width is not None else 0
        new_height = shape.height / 914400 if shape.height is not None else 0
        
        result = {
            "modified": len(changes_made) > 0,
            "element_index": element_index,
            "element_type": element_type,
            "changes_made": changes_made,
            "new_position": {
                "left": round(new_left, 3),
                "top": round(new_top, 3),
                "width": round(new_width, 3),
                "height": round(new_height, 3),
                "right": round(new_left + new_width, 3),
                "bottom": round(new_top + new_height, 3)
            }
        }

        return await save_and_return_presentation(result, presentation_id, context)

@slide_maker.action("get_element_styling")
class GetElementStylingAction(ActionHandler):
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        element_index = inputs["element_index"]
        files = inputs.get("files", [])

        load_presentation_from_files(presentation_id, files)

        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")

        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")

        slide = prs.slides[slide_index]
        if element_index >= len(slide.shapes):
            if len(slide.shapes) == 0:
                raise ValueError(f"Element index {element_index} out of range. Slide has no elements.")
            else:
                raise ValueError(f"Element index {element_index} out of range. Valid range: 0-{len(slide.shapes)-1} ({len(slide.shapes)} elements total).")

        shape = slide.shapes[element_index]

        # Get element type
        element_type = "unknown"
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type
            if shape_type == 1:
                element_type = "shape"
            elif shape_type == 13:
                element_type = "image"
            elif shape_type == 19:
                element_type = "table"
            elif shape_type == 3:
                element_type = "chart"
            elif shape_type == 17:
                element_type = "text_box"
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                element_type = "text_element"

        # Get position info
        left_inches = shape.left / 914400 if shape.left is not None else 0
        top_inches = shape.top / 914400 if shape.top is not None else 0
        width_inches = shape.width / 914400 if shape.width is not None else 0
        height_inches = shape.height / 914400 if shape.height is not None else 0

        styling_info = f"{element_type.upper()} @({left_inches:.1f},{top_inches:.1f}) {width_inches:.1f}×{height_inches:.1f}\":"

        if element_type == "table" and hasattr(shape, 'table'):
            styling_info += self._analyze_table_styling(shape.table)
        elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            styling_info += self._analyze_text_styling(shape)
        elif element_type == "image":
            styling_info += self._analyze_image_styling(shape)
        else:
            styling_info += "\nBASIC SHAPE: Limited styling information available"

        return {
            "element_index": element_index,
            "element_type": element_type,
            "styling_description": styling_info,
            "position": {
                "left": round(left_inches, 3),
                "top": round(top_inches, 3),
                "width": round(width_inches, 3),
                "height": round(height_inches, 3)
            }
        }

    def _analyze_table_styling(self, table):
        """Analyze table styling in compressed format"""
        rows = len(table.rows)
        cols = len(table.columns)

        info = f"\n\nTABLE {rows}×{cols}:\n"

        # Analyze font patterns across table
        font_patterns = self._analyze_table_fonts(table, rows, cols)

        # Build compressed description
        info += "PATTERNS:\n"
        info += f"- Structure: {rows} rows × {cols} columns\n"
        info += "- Borders: Standard table borders\n"

        # Add font information
        if font_patterns:
            info += f"- Fonts: {font_patterns}\n"

        # Check for header patterns
        patterns = []
        if rows > 0:
            first_row_bold = any(self._is_cell_bold(table.cell(0, c)) for c in range(cols))
            if first_row_bold:
                patterns.append("R0=headers(bold)")

            # Check first column for headers
            if cols > 0:
                first_col_bold = any(self._is_cell_bold(table.cell(r, 0)) for r in range(1, min(rows, 4)))
                if first_col_bold:
                    patterns.append("C0=row-labels(bold)")

        if patterns:
            info += "- Header patterns: " + ", ".join(patterns) + "\n"

        # Add column info with font details from data rows
        info += "\nCOLS:\n"
        col_descriptions = []
        for c in range(min(cols, 6)):  # Limit to 6 columns to save tokens
            # Try to get font info from data row (row 1) if available, otherwise header row
            sample_row = 1 if rows > 1 else 0
            sample_cell = table.cell(sample_row, c)
            sample_text = sample_cell.text.strip()[:20]
            font_info = self._get_cell_font_info(sample_cell)

            col_desc = f"C{c}: \"{sample_text}{'...' if len(sample_text) >= 20 else ''}\""
            if font_info:
                col_desc += f" ({font_info})"
            col_descriptions.append(col_desc)

        info += "- " + " | ".join(col_descriptions) + "\n"

        info += "\nRECREATE: add_table + modify_element for styling"

        return info

    def _analyze_table_fonts(self, table, rows, cols):
        """Analyze font patterns across the table"""
        font_samples = []

        # Expanded sampling to cover more cells including data rows
        sample_positions = [
            (0, 0),  # Header top-left
            (0, 1) if cols > 1 else None,  # Header second column
            (0, 2) if cols > 2 else None,  # Header third column
            (1, 0) if rows > 1 else None,  # First data row, first column
            (1, 1) if rows > 1 and cols > 1 else None,  # First data row, second column
            (1, 2) if rows > 1 and cols > 2 else None,  # First data row, third column (salary column)
        ]

        for pos in sample_positions:
            if pos:
                r, c = pos
                if r < rows and c < cols:
                    cell = table.cell(r, c)
                    font_info = self._get_cell_font_info(cell)
                    if font_info:
                        font_samples.append(f"({r},{c})={font_info}")

        if font_samples:
            return " ".join(font_samples[:6])  # Show up to 6 samples
        return "standard"

    def _get_cell_font_info(self, cell):
        """Extract font information from a table cell"""
        try:
            if not (cell.text_frame and cell.text_frame.paragraphs):
                return ""

            p = cell.text_frame.paragraphs[0]
            font_parts = []

            # Check paragraph font properties
            if hasattr(p, 'font') and p.font is not None:
                font = p.font

                # Font name
                if hasattr(font, 'name') and font.name:
                    font_parts.append(font.name)

                # Font size
                if hasattr(font, 'size') and font.size is not None:
                    try:
                        font_parts.append(f"{font.size.pt:.0f}pt")
                    except (AttributeError, TypeError):
                        pass

                # Font weight and style
                if hasattr(font, 'bold') and font.bold is not None and font.bold:
                    font_parts.append("bold")
                if hasattr(font, 'italic') and font.italic is not None and font.italic:
                    font_parts.append("italic")

                # Font color
                if hasattr(font, 'color') and font.color:
                    try:
                        if hasattr(font.color, 'rgb') and font.color.rgb:
                            rgb = font.color.rgb
                            if rgb != (0, 0, 0):  # Not black
                                if rgb == (255, 255, 255):
                                    font_parts.append("white")
                                else:
                                    hex_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                    font_parts.append(hex_color)
                    except:
                        pass

            # Check runs if paragraph font didn't work OR if we're missing font name
            if (not font_parts) or (font_parts and not any('Lato' in part or 'Arial' in part or 'Calibri' in part for part in font_parts)):
                # If we have some formatting but no font name, try runs for font name
                if p.runs:
                    for run in p.runs[:2]:  # Check first 2 runs
                        if hasattr(run, 'font') and run.font is not None:
                            run_font = run.font

                            # Font name from run - prioritize this
                            if hasattr(run_font, 'name') and run_font.name:
                                # Insert font name at the beginning if we don't have one
                                if not any(part for part in font_parts if not (part.endswith('pt') or part in ['bold', 'italic', 'white'] or part.startswith('#'))):
                                    font_parts.insert(0, run_font.name)
                                elif not font_parts:  # No font parts at all
                                    font_parts.append(run_font.name)

                            # Font size from run - only if we don't have it
                            if hasattr(run_font, 'size') and run_font.size is not None and not any('pt' in part for part in font_parts):
                                try:
                                    font_parts.append(f"{run_font.size.pt:.0f}pt")
                                except (AttributeError, TypeError):
                                    pass

                            # Bold/italic from run - only if we don't have it
                            if hasattr(run_font, 'bold') and run_font.bold is not None and run_font.bold and 'bold' not in font_parts:
                                font_parts.append("bold")
                            if hasattr(run_font, 'italic') and run_font.italic is not None and run_font.italic and 'italic' not in font_parts:
                                font_parts.append("italic")

                            # Color from run - only if we don't have it
                            if hasattr(run_font, 'color') and run_font.color and not any(part == 'white' or part.startswith('#') for part in font_parts):
                                try:
                                    if hasattr(run_font.color, 'rgb') and run_font.color.rgb:
                                        rgb = run_font.color.rgb
                                        if rgb != (0, 0, 0):  # Not black
                                            if rgb == (255, 255, 255):
                                                font_parts.append("white")
                                            else:
                                                hex_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                                font_parts.append(hex_color)
                                except:
                                    pass

                            if font_parts:  # Found something, we can stop
                                break

            return " ".join(font_parts) if font_parts else ""

        except Exception:
            return ""

    def _analyze_text_styling(self, shape):
        """Analyze text element styling"""
        info = "\n\nTEXT STYLING:\n"

        try:
            text_frame = shape.text_frame
            sample_text = text_frame.text[:50]
            if len(text_frame.text) > 50:
                sample_text += "..."

            info += f"- Content: \"{sample_text}\"\n"

            # Get comprehensive font info
            font_details = self._extract_text_font_info(text_frame)
            if font_details:
                info += f"- Font: {font_details}\n"

            # Get alignment info
            if text_frame.paragraphs:
                p = text_frame.paragraphs[0]
                if hasattr(p, 'alignment'):
                    align_map = {1: "center", 2: "right", 3: "justify"}
                    align = align_map.get(p.alignment, "left")
                    info += f"- Align: {align}\n"

            # Auto-sizing
            if hasattr(text_frame, 'auto_size'):
                auto_size_map = {0: "none", 1: "shape-to-text", 2: "text-to-shape"}
                auto_size = auto_size_map.get(text_frame.auto_size, "unknown")
                info += f"- Auto-size: {auto_size}\n"

            # Margins
            margin_info = []
            if hasattr(text_frame, 'margin_left') and text_frame.margin_left:
                margin_info.append(f"L{text_frame.margin_left / 914400:.1f}\"")
            if hasattr(text_frame, 'margin_right') and text_frame.margin_right:
                margin_info.append(f"R{text_frame.margin_right / 914400:.1f}\"")
            if hasattr(text_frame, 'margin_top') and text_frame.margin_top:
                margin_info.append(f"T{text_frame.margin_top / 914400:.1f}\"")
            if hasattr(text_frame, 'margin_bottom') and text_frame.margin_bottom:
                margin_info.append(f"B{text_frame.margin_bottom / 914400:.1f}\"")

            if margin_info:
                info += f"- Margins: {' '.join(margin_info)}\n"

        except Exception as e:
            info += f"- Error analyzing text: {str(e)[:50]}\n"

        info += "\nRECREATE: add_text with matching formatting"
        return info

    def _extract_text_font_info(self, text_frame):
        """Extract comprehensive font information from text frame"""
        try:
            if not text_frame.paragraphs:
                return ""

            font_parts = []

            # Try to get font from first paragraph
            p = text_frame.paragraphs[0]

            # Check paragraph font
            if hasattr(p, 'font') and p.font is not None:
                font = p.font

                # Font name
                if hasattr(font, 'name') and font.name:
                    font_parts.append(font.name)

                # Font size
                if hasattr(font, 'size') and font.size is not None:
                    try:
                        font_parts.append(f"{font.size.pt:.0f}pt")
                    except:
                        pass

                # Font weight and style
                if hasattr(font, 'bold') and font.bold is not None and font.bold:
                    font_parts.append("bold")
                if hasattr(font, 'italic') and font.italic is not None and font.italic:
                    font_parts.append("italic")

                # Font color
                if hasattr(font, 'color') and font.color:
                    try:
                        if hasattr(font.color, 'rgb') and font.color.rgb:
                            rgb = font.color.rgb
                            if rgb != (0, 0, 0):  # Not default black
                                if rgb == (255, 255, 255):
                                    font_parts.append("white")
                                else:
                                    hex_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                    font_parts.append(hex_color)
                    except:
                        pass

            # If no paragraph font info, check runs
            if not font_parts and p.runs:
                for run in p.runs[:2]:  # Check first 2 runs
                    if hasattr(run, 'font') and run.font is not None:
                        run_font = run.font

                        # Font name from run
                        if hasattr(run_font, 'name') and run_font.name:
                            font_parts.append(run_font.name)

                        # Font size from run
                        if hasattr(run_font, 'size') and run_font.size is not None:
                            try:
                                font_parts.append(f"{run_font.size.pt:.0f}pt")
                            except:
                                pass

                        # Bold/italic from run
                        if hasattr(run_font, 'bold') and run_font.bold is not None and run_font.bold:
                            font_parts.append("bold")
                        if hasattr(run_font, 'italic') and run_font.italic is not None and run_font.italic:
                            font_parts.append("italic")

                        if font_parts:  # If we found something, stop looking
                            break

            return " ".join(font_parts) if font_parts else ""

        except Exception:
            return ""

    def _analyze_image_styling(self, shape):
        """Analyze image styling"""
        info = "\n\nIMAGE STYLING:\n"
        info += "- Type: Picture/Image\n"
        info += "- Aspect: Preserved from original\n"
        info += "\nRECREATE: add_image with same position/size"
        return info

    def _get_cell_style_info(self, cell, row, col):
        """Get basic cell styling info"""
        try:
            text = cell.text.strip()
            is_bold = self._is_cell_bold(cell)
            return {
                "row": row, "col": col, "text": text[:20],
                "bold": is_bold
            }
        except:
            return {"row": row, "col": col, "text": "", "bold": False}

    def _is_cell_bold(self, cell):
        """Check if cell text is bold"""
        try:
            if cell.text_frame and cell.text_frame.paragraphs:
                for p in cell.text_frame.paragraphs:
                    if hasattr(p.font, 'bold') and p.font.bold:
                        return True
            return False
        except:
            return False

@slide_maker.action("find_and_replace")
class FindAndReplaceAction(ActionHandler):
    """Find and replace text across entire presentation with safety checks"""
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        replacements = inputs["replacements"]
        files = inputs.get("files", [])

        load_presentation_from_files(presentation_id, files)

        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")

        prs = presentations[presentation_id]
        total_replacements = 0
        warnings = []
        blocked_replacements = []
        changes = []  # Track detailed changelog for each replacement

        for replacement in replacements:
            find_text = replacement["find"]
            replace_text = replacement.get("replace", "")
            replace_all = replacement.get("replace_all", False)
            forced_font_size = replacement.get("forced_font_size")  # Optional: force specific size

            # Extract metadata from placeholder (e.g., "[Title, Fontsize=32pt, Bold=true]")
            _, placeholder_meta_dict = detect_placeholders_with_metadata(find_text)
            placeholder_meta = placeholder_meta_dict.get(find_text, {}) if placeholder_meta_dict else {}

            # Always strip ALL markdown from replacement text (form-filling doesn't parse markdown)
            # This prevents literal **/** appearing in output regardless of metadata
            replace_text = replace_text.replace('**', '').replace('__', '').replace('*', '').replace('_', '').replace('`', '').replace('~~', '')

            # Validation
            if not find_text or len(find_text.strip()) == 0:
                warnings.append(f"Skipped: 'find' text cannot be empty")
                continue

            # STEP 1: Scan entire presentation for matches
            matches_found = []

            for slide_idx, slide in enumerate(prs.slides):
                # Scan all text boxes
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        if find_text in shape.text_frame.text:
                            matches_found.append({
                                "type": "text_box",
                                "slide_index": slide_idx,
                                "element_index": shape_idx,
                                "content": shape.text_frame.text[:80] + "..." if len(shape.text_frame.text) > 80 else shape.text_frame.text,
                                "location": f"Slide {slide_idx}, Element {shape_idx}"
                            })

                # Scan all table cells
                for shape_idx, shape in enumerate(slide.shapes):
                    # Check if shape is actually a table (shape_type == 19)
                    if hasattr(shape, 'shape_type') and shape.shape_type == 19:
                        try:
                            table = shape.table
                            for row_idx, row in enumerate(table.rows):
                                for col_idx, cell in enumerate(row.cells):
                                    if find_text in cell.text:
                                        matches_found.append({
                                            "type": "table_cell",
                                            "slide_index": slide_idx,
                                            "element_index": shape_idx,
                                            "row": row_idx,
                                            "col": col_idx,
                                            "content": cell.text[:80] + "..." if len(cell.text) > 80 else cell.text,
                                            "location": f"Slide {slide_idx}, Table {shape_idx}, Cell ({row_idx},{col_idx})"
                                        })
                        except:
                            # Skip shapes that can't be accessed as tables
                            pass

            # STEP 2: Safety check
            if len(matches_found) == 0:
                warnings.append(f"No matches found for '{find_text}'")
                # Track as not found
                changes.append({
                    "find": find_text,
                    "replace": replace_text,
                    "status": "not_found",
                    "occurrences": 0,
                    "suggestion": "Placeholder not found in template - check spelling or verify template has this placeholder"
                })
                continue

            if len(matches_found) > 1 and not replace_all:
                # BLOCK! Multiple matches without explicit confirmation
                blocked_replacements.append({
                    "BLOCKED": f"'{find_text}' found {len(matches_found)} times - requires replace_all=true",
                    "find_phrase": find_text,
                    "match_count": len(matches_found),
                    "matches": matches_found[:5],  # Show first 5 matches
                    "fix_required": "Either add more context to find phrase to make it unique, OR set replace_all=true to confirm replacing all instances"
                })
                warnings.append(f"BLOCKED '{find_text}': {len(matches_found)} matches found, replace_all not set")
                # Track as blocked
                changes.append({
                    "find": find_text,
                    "replace": replace_text,
                    "status": "blocked",
                    "occurrences": len(matches_found),
                    "reason": "Multiple matches require replace_all=true confirmation"
                })
                continue

            # STEP 3: Safe to proceed - perform replacements
            replacement_count = 0

            # Initialize change record for successful replacement
            change_record = {
                "find": find_text,
                "replace": replace_text,
                "status": "replaced",
                "occurrences": len(matches_found),
                "locations": [],
                "font_size_applied": None
            }
            if forced_font_size:
                change_record["forced"] = True

            for match in matches_found:
                slide = prs.slides[match["slide_index"]]

                if match["type"] == "text_box":
                    # Replace in text box with markdown + auto-sizing
                    shape = slide.shapes[match["element_index"]]
                    text_frame = shape.text_frame

                    # FORM-FILLING APPROACH: Preserve all formatting, just replace text and resize

                    # Step 1: Capture original font size (template designer's intent)
                    original_font_size = 18  # Default fallback
                    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                        first_run = text_frame.paragraphs[0].runs[0]
                        if first_run.font.size:
                            original_font_size = int(first_run.font.size.pt)

                    # Step 2: In-place text replacement (preserves bullets, levels, colors, bold/italic)
                    # Handle both single-run and multi-run placeholders
                    for paragraph in text_frame.paragraphs:
                        # Try run-level replacement first (preserves run formatting perfectly)
                        replaced_in_run = False
                        for run in paragraph.runs:
                            if find_text in run.text:
                                run.text = run.text.replace(find_text, replace_text)
                                replaced_in_run = True

                        # If not found in any individual run, check if spans multiple runs
                        if not replaced_in_run and find_text in paragraph.text:
                            # Placeholder spans multiple runs - use paragraph-level replacement
                            new_text = paragraph.text.replace(find_text, replace_text)

                            # Clear all runs and set new text (keeps paragraph formatting)
                            for run in paragraph.runs:
                                run.text = ""

                            # Set first run to new text (or create if none)
                            if paragraph.runs:
                                paragraph.runs[0].text = new_text
                            else:
                                paragraph.text = new_text

                    # Step 3: Calculate size for ENTIRE text box (after replacement)
                    final_text = text_frame.text
                    width_inches = shape.width / 914400 if shape.width else 5
                    height_inches = shape.height / 914400 if shape.height else 1

                    if FONT_SIZE_DEBUG:
                        print(f"\n[FIND_REPLACE] Text box replacement:")
                        print(f"  Original placeholder size: {original_font_size}pt")
                        print(f"  Final text length: {len(final_text)} chars")
                        if forced_font_size:
                            print(f"  Forced font size requested: {forced_font_size}pt")
                        if placeholder_meta:
                            print(f"  Placeholder metadata: {placeholder_meta}")

                    # Determine font size (priority: forced > metadata > auto-calculate)
                    if forced_font_size:
                        # Priority 1: Explicit forced_font_size parameter
                        best_fit_size = forced_font_size
                        if FONT_SIZE_DEBUG:
                            print(f"  [FORCED SIZE] Using forced size: {forced_font_size}pt (skipping calculation)")
                    elif 'fontsize' in placeholder_meta:
                        # Priority 2: Fontsize in placeholder metadata
                        meta_size_str = placeholder_meta['fontsize']
                        # Strip 'pt' suffix (case-insensitive: 32pt, 32PT, 32Pt all work)
                        meta_size_str = meta_size_str.lower().replace('pt', '').strip()
                        try:
                            best_fit_size = int(meta_size_str)
                            if FONT_SIZE_DEBUG:
                                print(f"  [METADATA SIZE] Using size from placeholder: {best_fit_size}pt")
                        except:
                            # Invalid fontsize format, fall back to auto-calculate
                            has_markdown = has_markdown_formatting(final_text)
                            has_bullets = any(bullet in final_text for bullet in ['•', '◦', '▪', '▫', '‣', '-', '*'])
                            best_fit_size = calculate_best_fit_font_size(final_text, width_inches, height_inches, max_font_size=original_font_size, has_formatting=has_markdown, is_bullets=has_bullets, font_face='Calibri', debug=FONT_SIZE_DEBUG)
                    else:
                        # Priority 3: Auto-calculate
                        has_markdown = has_markdown_formatting(final_text)
                        has_bullets = any(bullet in final_text for bullet in ['•', '◦', '▪', '▫', '‣', '-', '*'])
                        best_fit_size = calculate_best_fit_font_size(final_text, width_inches, height_inches, max_font_size=original_font_size, has_formatting=has_markdown, is_bullets=has_bullets, font_face='Calibri', debug=FONT_SIZE_DEBUG)

                    # Step 4: Apply size and metadata-specified formatting uniformly to ALL runs
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Apply font size
                            run.font.size = Pt(best_fit_size)

                            # Apply metadata formatting if specified
                            if placeholder_meta:
                                # Bold
                                if 'bold' in placeholder_meta:
                                    run.font.bold = placeholder_meta['bold'].lower() == 'true'

                                # Italic
                                if 'italic' in placeholder_meta:
                                    run.font.italic = placeholder_meta['italic'].lower() == 'true'

                                # Underline
                                if 'underline' in placeholder_meta:
                                    run.font.underline = placeholder_meta['underline'].lower() == 'true'

                                # Font face
                                if 'font' in placeholder_meta:
                                    run.font.name = placeholder_meta['font']

                                # Color
                                if 'color' in placeholder_meta:
                                    color_value = placeholder_meta['color']
                                    if color_value.startswith('#'):
                                        # Hex color
                                        hex_color = color_value.lstrip('#')
                                        r = int(hex_color[0:2], 16)
                                        g = int(hex_color[2:4], 16)
                                        b = int(hex_color[4:6], 16)
                                        run.font.color.rgb = RGBColor(r, g, b)

                    # Disable auto-sizing to respect our calculated font size
                    text_frame.auto_size = MSO_AUTO_SIZE.NONE
                    text_frame.word_wrap = True

                    # Track this replacement location and font size
                    change_record["locations"].append({
                        "slide": match["slide_index"],
                        "element": match["element_index"],
                        "type": "text_box"
                    })
                    change_record["font_size_applied"] = best_fit_size

                    replacement_count += 1

                elif match["type"] == "table_cell":
                    # Replace in table cell with markdown + auto-sizing
                    try:
                        shape = slide.shapes[match["element_index"]]
                        table = shape.table
                        cell = table.cell(match["row"], match["col"])

                        # FORM-FILLING APPROACH: Preserve all formatting, just replace text and resize

                        # Step 1: Capture original font size
                        original_cell_font_size = 14  # Default fallback
                        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                            first_run = cell.text_frame.paragraphs[0].runs[0]
                            if first_run.font.size:
                                original_cell_font_size = int(first_run.font.size.pt)

                        # Step 2: In-place text replacement (preserves structure)
                        # Handle both single-run and multi-run placeholders
                        for paragraph in cell.text_frame.paragraphs:
                            # Try run-level replacement first
                            replaced_in_run = False
                            for run in paragraph.runs:
                                if find_text in run.text:
                                    run.text = run.text.replace(find_text, replace_text)
                                    replaced_in_run = True

                            # If not found in any individual run, check if spans multiple runs
                            if not replaced_in_run and find_text in paragraph.text:
                                # Placeholder spans multiple runs - use paragraph-level replacement
                                new_text = paragraph.text.replace(find_text, replace_text)

                                # Clear all runs and set new text
                                for run in paragraph.runs:
                                    run.text = ""

                                # Set first run to new text
                                if paragraph.runs:
                                    paragraph.runs[0].text = new_text
                                else:
                                    paragraph.text = new_text

                        # Step 3: Calculate size for entire cell (after replacement)
                        final_text = cell.text
                        cell_width = table.columns[match["col"]].width / 914400
                        cell_height = table.rows[match["row"]].height / 914400

                        # Determine font size (priority: forced > metadata > auto-calculate)
                        if forced_font_size:
                            best_fit_size = forced_font_size
                            if FONT_SIZE_DEBUG:
                                print(f"  [FORCED SIZE] Table cell using forced size: {forced_font_size}pt")
                        elif 'fontsize' in placeholder_meta:
                            meta_size_str = placeholder_meta['fontsize']
                            # Strip 'pt' suffix (case-insensitive)
                            meta_size_str = meta_size_str.lower().replace('pt', '').strip()
                            try:
                                best_fit_size = int(meta_size_str)
                                if FONT_SIZE_DEBUG:
                                    print(f"  [METADATA SIZE] Table cell using size from placeholder: {best_fit_size}pt")
                            except:
                                has_markdown = has_markdown_formatting(final_text)
                                has_bullets = any(bullet in final_text for bullet in ['•', '◦', '▪', '▫', '‣', '-', '*'])
                                best_fit_size = calculate_best_fit_font_size(final_text, cell_width, cell_height, max_font_size=original_cell_font_size, has_formatting=has_markdown, is_bullets=has_bullets, font_face='Calibri', debug=FONT_SIZE_DEBUG)
                        else:
                            has_markdown = has_markdown_formatting(final_text)
                            has_bullets = any(bullet in final_text for bullet in ['•', '◦', '▪', '▫', '‣', '-', '*'])
                            best_fit_size = calculate_best_fit_font_size(final_text, cell_width, cell_height, max_font_size=original_cell_font_size, has_formatting=has_markdown, is_bullets=has_bullets, font_face='Calibri', debug=FONT_SIZE_DEBUG)

                        # Step 4: Apply size and metadata-specified formatting uniformly to all runs
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(best_fit_size)

                                # Apply metadata formatting if specified
                                if placeholder_meta:
                                    if 'bold' in placeholder_meta:
                                        run.font.bold = placeholder_meta['bold'].lower() == 'true'
                                    if 'italic' in placeholder_meta:
                                        run.font.italic = placeholder_meta['italic'].lower() == 'true'
                                    if 'underline' in placeholder_meta:
                                        run.font.underline = placeholder_meta['underline'].lower() == 'true'
                                    if 'font' in placeholder_meta:
                                        run.font.name = placeholder_meta['font']
                                    if 'color' in placeholder_meta:
                                        color_value = placeholder_meta['color']
                                        if color_value.startswith('#'):
                                            hex_color = color_value.lstrip('#')
                                            r = int(hex_color[0:2], 16)
                                            g = int(hex_color[2:4], 16)
                                            b = int(hex_color[4:6], 16)
                                            run.font.color.rgb = RGBColor(r, g, b)

                        cell.text_frame.word_wrap = True

                        # Track this replacement location and font size
                        change_record["locations"].append({
                            "slide": match["slide_index"],
                            "element": match["element_index"],
                            "type": "table_cell",
                            "row": match["row"],
                            "col": match["col"]
                        })
                        change_record["font_size_applied"] = best_fit_size

                        replacement_count += 1
                    except:
                        # Skip if table access fails
                        pass

            # Add this replacement's change record to changes array
            if change_record.get("locations"):
                changes.append(change_record)

            total_replacements += replacement_count

        # Calculate summary statistics
        successful_count = sum(1 for c in changes if c["status"] == "replaced")
        failed_count = sum(1 for c in changes if c["status"] == "not_found")
        blocked_count = sum(1 for c in changes if c["status"] == "blocked")

        # Determine overall status
        if successful_count == len(replacements):
            status = "all_successful"
        elif blocked_count == len(replacements):
            status = "all_blocked"
        elif failed_count == len(replacements):
            status = "all_failed"
        elif successful_count > 0:
            status = "partial_success"
        else:
            status = "all_failed"

        result = {
            # New clear fields
            "status": status,
            "summary": {
                "requested": len(replacements),
                "successful": successful_count,
                "failed": failed_count,
                "blocked": blocked_count
            },
            "changes": changes,

            # Existing fields (kept for backward compatibility)
            "success": total_replacements > 0 and len(blocked_replacements) == 0,
            "total_replacements": total_replacements,
            "processed": len(replacements),
            "blocked": blocked_replacements,
            "warnings": warnings
        }

        return await save_and_return_presentation(result, presentation_id, context)


@slide_maker.action("add_elements")
class AddElementsAction(ActionHandler):
    """
    Unified action to add one or multiple elements from markdown content.
    Supports TWO MODES:
    1. GRANULAR MODE (auto_layout=false): Array of elements with position control
    2. AUTO-LAYOUT MODE (auto_layout=true): Single markdown document with vertical flow
    """
    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext):
        presentation_id = inputs["presentation_id"]
        slide_index = inputs["slide_index"]
        auto_layout = inputs.get("auto_layout", False)
        files = inputs.get("files", [])

        load_presentation_from_files(presentation_id, files)

        if presentation_id not in presentations:
            raise ValueError(f"Presentation {presentation_id} not found")

        prs = presentations[presentation_id]
        if slide_index >= len(prs.slides):
            if len(prs.slides) == 0:
                raise ValueError(f"Slide index {slide_index} out of range. Presentation has no slides.")
            else:
                raise ValueError(f"Slide index {slide_index} out of range. Valid range: 0-{len(prs.slides)-1} ({len(prs.slides)} slides total).")

        slide = prs.slides[slide_index]
        slide_width = prs.slide_width / 914400
        slide_height = prs.slide_height / 914400
        slide_dims = {"width": slide_width, "height": slide_height}

        # MODE DETECTION: Auto-layout or Granular
        if auto_layout:
            # AUTO-LAYOUT MODE: Single markdown document → vertical flow
            markdown_content = inputs.get("markdown")
            if not markdown_content:
                raise ValueError("auto_layout=true requires 'markdown' parameter")

            # Use existing parse_markdown_to_slide function (supports h1-h6, blockquotes, code, etc.)
            created_shapes = parse_markdown_to_slide(slide, markdown_content, prs)

            result = {
                "mode": "auto_layout",
                "total_requested": len(created_shapes),
                "successfully_added": len(created_shapes),
                "skipped": 0,
                "elements_added": [
                    {
                        "index": idx,
                        "type": shape_info["type"],
                        "shape_id": shape_info["shape_id"],
                        "requested_position": None,  # Auto-positioned
                        "final_position": None,  # Position info not tracked in auto-layout
                        "position_adjusted": False,
                        "adjustment_reason": "Auto-positioned by layout algorithm"
                    }
                    for idx, shape_info in enumerate(created_shapes)
                ],
                "elements_skipped": []
            }

            return await save_and_return_presentation(result, presentation_id, context)

        # GRANULAR MODE: Array of elements with position control
        elements_to_add = inputs.get("elements")
        if not elements_to_add:
            raise ValueError("auto_layout=false requires 'elements' parameter")

        # Get current elements for overlap detection
        existing_elements = get_current_slide_elements(slide, slide_width, slide_height)

        elements_added = []
        elements_skipped = []

        for idx, element in enumerate(elements_to_add):
            content = element["content"]
            requested_pos = element.get("position")
            auto_position = element.get("auto_position", False)

            # STEP 1: Auto-detect element type from markdown
            element_type = detect_element_type_from_markdown(content)

            # STEP 2: Parse markdown content based on type
            try:
                if element_type == "table":
                    parsed_data = parse_markdown_table(content)
                    if parsed_data["rows"] == 0 or parsed_data["cols"] == 0:
                        elements_skipped.append({
                            "index": idx,
                            "content_preview": content[:50] + "..." if len(content) > 50 else content,
                            "skip_reason": "Failed to parse markdown table - check table syntax",
                            "suggestion": "Ensure table has | separators and a header separator line (|---|---|)"
                        })
                        continue
                    # Set default dimensions if not provided
                    if not requested_pos:
                        requested_pos = {
                            "left": 0.5,
                            "top": 0.5 + (idx * 0.5),
                            "width": 5.0,
                            "height": 1.0 + (parsed_data["rows"] * 0.3)
                        }
                elif element_type == "bullets":
                    parsed_data = parse_markdown_bullets(content)
                    if len(parsed_data) == 0:
                        elements_skipped.append({
                            "index": idx,
                            "content_preview": content[:50] + "..." if len(content) > 50 else content,
                            "skip_reason": "Failed to parse markdown bullets - check list syntax",
                            "suggestion": "Use - or * for unordered lists, or 1. 2. for ordered lists"
                        })
                        continue
                    if not requested_pos:
                        requested_pos = {
                            "left": 0.5,
                            "top": 0.5 + (idx * 0.5),
                            "width": 4.0,
                            "height": 1.5 + (len(parsed_data) * 0.2)
                        }
                else:  # text
                    parsed_data = content
                    if not requested_pos:
                        requested_pos = {
                            "left": 0.5,
                            "top": 0.5 + (idx * 0.5),
                            "width": 4.0,
                            "height": 1.0
                        }

                # STEP 3: Validate and adjust position
                validation = validate_and_adjust_position(
                    requested_pos, existing_elements, slide_dims, auto_position
                )

                if not validation["fits"]:
                    elements_skipped.append({
                        "index": idx,
                        "content_preview": content[:50] + "..." if len(content) > 50 else content,
                        "skip_reason": validation["reason"],
                        "suggestion": "Try enabling auto_position, providing different coordinates, or adjusting element size"
                    })
                    continue

                final_pos = validation["position"]

                # STEP 4: Create element using internal helpers
                if element_type == "text":
                    shape = _create_text_box(slide, parsed_data, final_pos, prs)
                elif element_type == "table":
                    shape = _create_table(slide, parsed_data, final_pos, prs)
                elif element_type == "bullets":
                    shape = _create_bullet_list(slide, parsed_data, final_pos, prs)

                elements_added.append({
                    "index": idx,
                    "type": element_type,
                    "shape_id": str(shape.shape_id),
                    "requested_position": requested_pos if element.get("position") else None,
                    "final_position": final_pos,
                    "position_adjusted": validation["adjusted"],
                    "adjustment_reason": validation.get("reason") or "No adjustment needed"
                })

                # Update existing_elements for next iteration
                existing_elements.append({
                    "index": len(existing_elements),
                    "position": {
                        **final_pos,
                        "right": final_pos["left"] + final_pos["width"],
                        "bottom": final_pos["top"] + final_pos["height"]
                    }
                })

            except Exception as e:
                elements_skipped.append({
                    "index": idx,
                    "content_preview": content[:50] + "..." if len(content) > 50 else content,
                    "skip_reason": f"Error creating element: {str(e)}",
                    "suggestion": "Check markdown syntax and ensure dimensions are valid"
                })

        result = {
            "mode": "granular",
            "total_requested": len(elements_to_add),
            "successfully_added": len(elements_added),
            "skipped": len(elements_skipped),
            "elements_added": elements_added,
            "elements_skipped": elements_skipped
        }

        return await save_and_return_presentation(result, presentation_id, context)



