# Test the get_slide_elements action locally to debug text overflow calculations
import asyncio
import os
import sys
from context import slide_maker
from autohive_integrations_sdk import ExecutionContext

async def load_and_inspect_presentation(presentation_path):
    """Load a presentation and inspect all slides"""
    print(f"Loading presentation from: {presentation_path}")
    
    # Check if file exists
    if not os.path.exists(presentation_path):
        raise FileNotFoundError(f"Presentation file not found: {presentation_path}")
    
    # Load the presentation directly into memory (for local testing)
    from slide_maker import presentations
    from pptx import Presentation
    import uuid
    
    # Load the presentation directly with python-pptx
    prs = Presentation(presentation_path)
    presentation_id = os.path.basename(presentation_path)  # Use filename as ID
    
    # Store it in the presentations dictionary
    presentations[presentation_id] = prs
    
    print(f"Loaded presentation with {len(prs.slides)} slides")
    print(f"Using presentation ID: {presentation_id}")
    
    # Mock auth for ExecutionContext
    auth = {}
    
    # Inspect each slide
    async with ExecutionContext(auth=auth) as context:
        for slide_idx in range(len(prs.slides)):
            print(f"\n" + "="*60)
            print(f"SLIDE {slide_idx} INSPECTION")
            print("="*60)
            
            # Run inspection with content enabled
            inspect_inputs = {
                "presentation_id": presentation_id,
                "slide_index": slide_idx,
                "include_content": True,  # Always show content for debugging
                "files": []  # No files needed for local testing
            }
            
            result = await slide_maker.execute_action("get_slide_elements", inspect_inputs, context)
            
            # Display results
            print(f"Slide {slide_idx} Summary:")
            print(f"  - Total elements: {result['total_elements']}")
            print(f"  - Layout status: {result['layout_status']}")
            print(f"  - Slide dimensions: {result['slide_dimensions']['width']}\" x {result['slide_dimensions']['height']}\"")
            
            if result['layout_status'] == 'issues_detected':
                if 'elements_outside_boundary' in result:
                    print(f"  - Elements outside boundary: {result['elements_outside_boundary']}")
                if 'elements_with_text_overflow' in result:
                    print(f"  - Elements with text overflow: {result['elements_with_text_overflow']}")
                if 'total_overlapping_pairs' in result:
                    print(f"  - Overlapping pairs: {result['total_overlapping_pairs']}")
            
            # Display each element
            print(f"\nElements on slide {slide_idx}:")
            for element in result['elements']:
                print(f"\n  Element {element['index']} ({element['type']}):")
                pos = element['position']
                print(f"    Position: {pos['left']}\", {pos['top']}\" | Size: {pos['width']}\" x {pos['height']}\"")
                
                if 'content' in element:
                    content_preview = element['content'][:100] + "..." if len(element['content']) > 100 else element['content']
                    print(f"    Content: \"{content_preview}\"")
                
                if 'debug_info' in element:
                    print(f"    Debug: {element['debug_info']}")
                
                if 'boundary_status' in element and element['boundary_status'] != 'inside':
                    print(f"    Boundary Status: {element['boundary_status']}")
                
                if 'boundary_warnings' in element:
                    for warning in element['boundary_warnings']:
                        print(f"    Warning: {warning}")
                
                if element.get('text_overflow_detected'):
                    print(f"    ⚠️  TEXT OVERFLOW DETECTED")
            
            # Display overlaps if any
            if 'element_overlaps' in result and result['element_overlaps']:
                print(f"\nOverlaps detected on slide {slide_idx}:")
                for overlap in result['element_overlaps']:
                    print(f"  - {overlap['description']}")
                    print(f"    Overlap area: {overlap['overlap_area']} sq inches ({overlap['severity']} severity)")

async def main():
    print("Local Presentation Inspection Tool")
    print("=================================")
    
    # Check if presentation file was provided as command line argument
    if len(sys.argv) < 2:
        print("Usage: python test_inspection.py <path_to_presentation.pptx>")
        print("Example: python test_inspection.py \"C:\\Users\\Onil\\Downloads\\bee_presentation.pptx\"")
        return
    
    presentation_path = sys.argv[1]
    
    try:
        await load_and_inspect_presentation(presentation_path)
        print(f"\n✅ Inspection completed successfully!")
        
    except Exception as e:
        print(f"\n❌ Test failed with error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    asyncio.run(main())