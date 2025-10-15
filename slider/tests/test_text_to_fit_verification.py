# Verify that TEXT_TO_FIT_SHAPE actually works when file is opened
import sys
import os
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "../dependencies")))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

print("Creating test to verify TEXT_TO_FIT_SHAPE behavior...")

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Box 1: Long text WITHOUT auto-sizing (control)
print("\n1. Box WITHOUT auto-sizing (control)...")
box1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
tf1 = box1.text_frame
tf1.word_wrap = True
tf1.margin_left = Inches(0.1)
tf1.margin_right = Inches(0.1)
p1 = tf1.paragraphs[0]
run1 = p1.add_run()
run1.text = "This is very long text that will overflow the box. " * 3
run1.font.name = 'Calibri'
run1.font.size = Pt(18)
print(f"   Font size: {run1.font.size.pt}pt")
print(f"   auto_size: {tf1.auto_size}")
print(f"   word_wrap: {tf1.word_wrap}")

# Box 2: Long text WITH TEXT_TO_FIT_SHAPE
print("\n2. Box WITH TEXT_TO_FIT_SHAPE...")
box2 = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(3), Inches(1))
tf2 = box2.text_frame
tf2.word_wrap = True
tf2.margin_left = Inches(0.1)
tf2.margin_right = Inches(0.1)
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = "This is very long text that will overflow the box. " * 3
run2.font.name = 'Calibri'
run2.font.size = Pt(18)
# Set auto-size AFTER adding text
tf2.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
print(f"   Font size: {run2.font.size.pt}pt")
print(f"   auto_size: {tf2.auto_size}")
print(f"   word_wrap: {tf2.word_wrap}")

# Box 3: Very long text WITH TEXT_TO_FIT_SHAPE
print("\n3. Box WITH TEXT_TO_FIT_SHAPE (extreme)...")
box3 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(0.8))
tf3 = box3.text_frame
tf3.word_wrap = True
tf3.margin_left = Inches(0.1)
tf3.margin_right = Inches(0.1)
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = "Extremely long text that definitely needs to shrink significantly. " * 5
run3.font.name = 'Calibri'
run3.font.size = Pt(18)
# Set auto-size AFTER adding text
tf3.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
print(f"   Font size: {run3.font.size.pt}pt")
print(f"   auto_size: {tf3.auto_size}")
print(f"   word_wrap: {tf3.word_wrap}")

# Box 4: Multiple runs WITH TEXT_TO_FIT_SHAPE (like markdown)
print("\n4. Box with MULTIPLE RUNS + TEXT_TO_FIT_SHAPE...")
box4 = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(2), Inches(0.8))
tf4 = box4.text_frame
tf4.word_wrap = True
tf4.margin_left = Inches(0.1)
tf4.margin_right = Inches(0.1)
p4 = tf4.paragraphs[0]

# Add multiple runs with formatting (simulating markdown)
r1 = p4.add_run()
r1.text = "Bold text "
r1.font.name = 'Calibri'
r1.font.size = Pt(18)
r1.font.bold = True

r2 = p4.add_run()
r2.text = "and italic text "
r2.font.name = 'Calibri'
r2.font.size = Pt(18)
r2.font.italic = True

r3 = p4.add_run()
r3.text = "with lots more content to overflow. " * 3
r3.font.name = 'Calibri'
r3.font.size = Pt(18)

# Set auto-size AFTER adding all runs
tf4.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
print(f"   Runs: 3 (bold, italic, normal)")
print(f"   All font sizes: 18pt")
print(f"   All font names: Calibri")
print(f"   auto_size: {tf4.auto_size}")
print(f"   word_wrap: {tf4.word_wrap}")

# Save
output_path = os.path.join(os.path.dirname(__file__), "..", "test_text_to_fit_verification.pptx")
prs.save(output_path)

print("\n" + "="*70)
print(f"SAVED: {output_path}")
print("="*70)
print("\nCRITICAL TEST:")
print("  Close this terminal, then open the PowerPoint file.")
print("\n  Check if the text AUTOMATICALLY SHRUNK when you opened it:")
print("    - Box 1 (top-left):    Should OVERFLOW (no auto-sizing)")
print("    - Box 2 (top-right):   Should SHRINK TO FIT automatically")
print("    - Box 3 (bottom-left): Should SHRINK TO FIT automatically")
print("    - Box 4 (bottom-right): Should SHRINK TO FIT (with formatting intact)")
print("\n  If boxes 2, 3, 4 do NOT shrink, TEXT_TO_FIT_SHAPE doesn't work!")
print("="*70)
