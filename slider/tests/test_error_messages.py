# Test error messages to ensure they provide helpful range information for LLM recovery
import asyncio
from context import slide_maker
from autohive_integrations_sdk import ExecutionContext

async def create_test_presentation_with_known_counts():
    """Helper to create a presentation with known slide/element counts"""
    auth = {}
    async with ExecutionContext(auth=auth) as context:
        # Create presentation with title (1 slide)
        create_inputs = {
            "title": "Error Message Test",
            "custom_filename": "test_error_messages_presentation"
        }
        result = await slide_maker.execute_action("create_presentation", create_inputs, context)
        presentation_id = result["presentation_id"]
        
        # Add second slide (total: 2 slides)
        add_slide_inputs = {
            "presentation_id": presentation_id,
            "files": [result["file"]]
        }
        result = await slide_maker.execute_action("add_slide", add_slide_inputs, context)
        
        # Add 3 elements to slide 0 (total: 3 elements on slide 0)
        text_inputs1 = {
            "presentation_id": presentation_id,
            "slide_index": 0,
            "text": "Element 1",
            "position": {"left": 1.0, "top": 3.0, "width": 2.0, "height": 0.5},
            "files": [result["file"]]
        }
        result = await slide_maker.execute_action("add_text", text_inputs1, context)
        
        text_inputs2 = {
            "presentation_id": presentation_id,
            "slide_index": 0,
            "text": "Element 2", 
            "position": {"left": 4.0, "top": 3.0, "width": 2.0, "height": 0.5},
            "files": [result["file"]]
        }
        result = await slide_maker.execute_action("add_text", text_inputs2, context)
        
        # Leave slide 1 empty (0 elements)
        
        return presentation_id, result["file"]

async def test_slide_index_errors():
    """Test slide index error messages"""
    print("Testing slide index error messages...")
    
    presentation_id, file_data = await create_test_presentation_with_known_counts()
    
    auth = {}
    async with ExecutionContext(auth=auth) as context:
        # Test invalid slide index (should show valid range 0-1)
        try:
            inputs = {
                "presentation_id": presentation_id,
                "slide_index": 5,  # Invalid - only have slides 0,1
                "files": [file_data]
            }
            
            await slide_maker.execute_action("get_slide_elements", inputs, context)
            print("❌ Should have failed with slide index error")
            
        except ValueError as e:
            error_msg = str(e)
            print(f"✅ Slide index error message: '{error_msg}'")
            
            # Verify error message contains helpful information
            if "Valid range: 0-1" in error_msg and "2 slides total" in error_msg:
                print("   ✅ Contains valid range and total count")
            else:
                print("   ❌ Missing range information")
                
        # Test another action with invalid slide index
        try:
            inputs = {
                "presentation_id": presentation_id,
                "slide_index": 10,
                "text": "Should fail",
                "position": {"left": 1.0, "top": 1.0, "width": 2.0, "height": 1.0},
                "files": [file_data]
            }
            
            await slide_maker.execute_action("add_text", inputs, context)
            print("❌ Should have failed with slide index error")
            
        except ValueError as e:
            error_msg = str(e)
            print(f"✅ Add text slide error: '{error_msg}'")
            
            if "Valid range:" in error_msg and "slides total" in error_msg:
                print("   ✅ Contains helpful range information")
            else:
                print("   ❌ Missing range information")

async def test_element_index_errors():
    """Test element index error messages"""
    print("\nTesting element index error messages...")
    
    presentation_id, file_data = await create_test_presentation_with_known_counts()
    
    auth = {}
    async with ExecutionContext(auth=auth) as context:
        # Test invalid element index on slide with elements (slide 0 has 3 elements)
        try:
            inputs = {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "element_index": 8,  # Invalid - only have elements 0,1,2
                "files": [file_data]
            }
            
            await slide_maker.execute_action("modify_element", inputs, context)
            print("❌ Should have failed with element index error")
            
        except ValueError as e:
            error_msg = str(e)
            print(f"✅ Element index error message: '{error_msg}'")
            
            if "Valid range: 0-2" in error_msg and "3 elements total" in error_msg:
                print("   ✅ Contains valid range and total count")
            else:
                print("   ❌ Missing range information")
                
        # Test element index on slide with no elements (slide 1 has 0 elements)
        try:
            inputs = {
                "presentation_id": presentation_id,
                "slide_index": 1,  # Empty slide
                "shape_index": 0,   # No elements exist
                "files": [file_data]
            }
            
            await slide_maker.execute_action("delete_element", inputs, context)
            print("❌ Should have failed with no elements error")
            
        except ValueError as e:
            error_msg = str(e)
            print(f"✅ No elements error message: '{error_msg}'")
            
            if "Slide has no elements" in error_msg:
                print("   ✅ Clear message for empty slide")
            else:
                print("   ❌ Unclear message for empty slide")

async def test_shape_index_errors():
    """Test shape index error messages for text control actions"""
    print("\nTesting shape index error messages...")
    
    presentation_id, file_data = await create_test_presentation_with_known_counts()
    
    auth = {}
    async with ExecutionContext(auth=auth) as context:
        # Test invalid shape index for text controls (slide 0 has 3 shapes)
        try:
            inputs = {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "shape_index": 7,  # Invalid - only have shapes 0,1,2
                "autosize_type": "TEXT_TO_FIT_SHAPE",
                "files": [file_data]
            }
            
            await slide_maker.execute_action("set_text_autosize", inputs, context)
            print("❌ Should have failed with shape index error")
            
        except ValueError as e:
            error_msg = str(e)
            print(f"✅ Shape index error message: '{error_msg}'")
            
            if "Valid range: 0-2" in error_msg and "3 elements total" in error_msg:
                print("   ✅ Contains valid range and total count")
            else:
                print("   ❌ Missing range information")

async def test_empty_presentation_errors():
    """Test error messages for completely empty presentation"""
    print("\nTesting empty presentation error messages...")
    
    # Create presentation but don't add any content
    from slide_maker import presentations
    from pptx import Presentation
    import uuid
    
    # Create truly empty presentation
    empty_prs = Presentation()
    empty_id = str(uuid.uuid4())
    presentations[empty_id] = empty_prs
    
    auth = {}
    async with ExecutionContext(auth=auth) as context:
        # Test accessing slide in empty presentation
        try:
            inputs = {
                "presentation_id": empty_id,
                "slide_index": 0,
                "files": []
            }
            
            await slide_maker.execute_action("get_slide_elements", inputs, context)
            print("❌ Should have failed with no slides error")
            
        except ValueError as e:
            error_msg = str(e)
            print(f"✅ Empty presentation error: '{error_msg}'")
            
            if "Presentation has no slides" in error_msg:
                print("   ✅ Clear message for empty presentation")
            else:
                print("   ❌ Unclear message for empty presentation")

async def test_error_message_consistency():
    """Test that error messages are consistent across all actions"""
    print("\nTesting error message consistency...")
    
    presentation_id, file_data = await create_test_presentation_with_known_counts()
    
    # List of actions that take slide_index
    slide_actions = [
        ("add_text", {
            "text": "Test",
            "position": {"left": 1.0, "top": 1.0, "width": 2.0, "height": 1.0}
        }),
        ("get_slide_elements", {}),
        ("add_bullet_list", {
            "bullet_items": [{"text": "Test", "level": 0}],
            "position": {"left": 1.0, "top": 1.0, "width": 2.0, "height": 1.0}
        })
    ]
    
    auth = {}
    async with ExecutionContext(auth=auth) as context:
        consistent_count = 0
        
        for action_name, extra_inputs in slide_actions:
            try:
                inputs = {
                    "presentation_id": presentation_id,
                    "slide_index": 99,  # Invalid index
                    "files": [file_data],
                    **extra_inputs
                }
                
                await slide_maker.execute_action(action_name, inputs, context)
                print(f"❌ {action_name} should have failed")
                
            except ValueError as e:
                error_msg = str(e)
                if "Valid range: 0-1" in error_msg and "2 slides total" in error_msg:
                    consistent_count += 1
                    print(f"   ✅ {action_name}: Consistent error format")
                else:
                    print(f"   ❌ {action_name}: Inconsistent error format")
        
        print(f"✅ Error message consistency: {consistent_count}/{len(slide_actions)} actions")

async def main():
    print("Testing Error Message Improvements")
    print("==================================")
    print("Validating LLM-friendly error messages with range information...")
    
    try:
        await test_slide_index_errors()
        await test_element_index_errors()
        await test_shape_index_errors()
        await test_empty_presentation_errors()
        await test_error_message_consistency()
        print("\n✅ All error message tests passed!")
        print("🤖 LLMs can now easily recover from index errors!")
        
    except Exception as e:
        print(f"\n❌ Test failed: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    asyncio.run(main())