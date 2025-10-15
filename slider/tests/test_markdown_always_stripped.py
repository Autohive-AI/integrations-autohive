"""
Test that markdown is ALWAYS stripped in find_and_replace,
even when placeholder has no metadata.
"""

import asyncio
import sys
import os
import base64

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_maker import slide_maker
from autohive_integrations_sdk import ExecutionContext
from pptx import Presentation
from io import BytesIO


async def test():
    print("=" * 70)
    print("MARKDOWN ALWAYS STRIPPED TEST")
    print("=" * 70)

    auth = {}

    async with ExecutionContext(auth=auth) as context:
        result = await slide_maker.execute_action(
            "create_presentation",
            {"title": "Markdown Strip Test"},
            context
        )

        presentation_id = result["presentation_id"]
        files = [result["file"]]

        # Add placeholders WITHOUT metadata
        await slide_maker.execute_action(
            "add_elements",
            {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "elements": [
                    {"content": "[Title]", "position": {"left": 1, "top": 1, "width": 8, "height": 1}},
                    {"content": "[Company]", "position": {"left": 1, "top": 2.2, "width": 8, "height": 0.8}},
                    {"content": "[Body]", "position": {"left": 1, "top": 3.2, "width": 8, "height": 1}}
                ],
                "files": files
            },
            context
        )

        print("\n1. Created placeholders (NO metadata)")

        # Agent sends markdown (should be stripped)
        print("\n2. Replacing with markdown (should be stripped)...")

        replacements = [
            {"find": "[Title]", "replace": "**Q4 Business Review**"},
            {"find": "[Company]", "replace": "*TechVenture Solutions*"},
            {"find": "[Body]", "replace": "This has `code` and ~~strikethrough~~ text"}
        ]

        result = await slide_maker.execute_action(
            "find_and_replace",
            {
                "presentation_id": presentation_id,
                "replacements": replacements,
                "files": files
            },
            context
        )

        print(f"   Status: {result['status']}")

        # Verify markdown was stripped
        prs_bytes = base64.b64decode(result["file"]["content"])
        prs = Presentation(BytesIO(prs_bytes))
        slide = prs.slides[0]

        print("\n3. Verification:")

        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text

                print(f"\n   Element {idx}: \"{text}\"")

                # Check for literal markdown characters
                has_literal_markdown = any(char in text for char in ['**', '*', '`', '~~', '__', '_'])

                if has_literal_markdown:
                    print(f"     [FAIL] Contains literal markdown characters!")
                    if '**' in text:
                        print(f"       Found: ** (should be stripped)")
                    if '`' in text:
                        print(f"       Found: ` (should be stripped)")
                else:
                    print(f"     [OK] No literal markdown characters")

        print("\n" + "=" * 70)
        print("EXPECTED:")
        print("  Element 0: 'Q4 Business Review' (no ** visible)")
        print("  Element 1: 'TechVenture Solutions' (no * visible)")
        print("  Element 2: 'This has code and strikethrough text' (no ` or ~~ visible)")
        print("=" * 70)


if __name__ == "__main__":
    asyncio.run(test())
