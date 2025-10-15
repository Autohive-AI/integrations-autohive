# Simple test to verify fit_text() actually works
import sys
import os
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "../dependencies")))

from pptx import Presentation
from pptx.util import Inches, Pt

print("Testing fit_text() functionality...")

# Create presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Test 1: Long text in small box WITHOUT fit_text
print("\n1. Creating text box WITHOUT fit_text()...")
box1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
tf1 = box1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
run1 = p1.add_run()
run1.text = "This is a very long text that will definitely overflow the small box if we don't shrink it"
run1.font.name = 'Calibri'
run1.font.size = Pt(18)
print(f"   Text length: {len(run1.text)} chars")
print(f"   Box size: 2x1 inches")
print(f"   Font size BEFORE: {run1.font.size.pt}pt")

# Test 2: Long text in small box WITH fit_text
print("\n2. Creating text box WITH fit_text()...")
box2 = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1))
tf2 = box2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = "This is a very long text that will definitely overflow the small box if we don't shrink it"
run2.font.name = 'Calibri'
run2.font.size = Pt(18)
print(f"   Text length: {len(run2.text)} chars")
print(f"   Box size: 2x1 inches")
print(f"   Font size BEFORE fit_text: {run2.font.size.pt}pt")

# Apply fit_text
try:
    tf2.fit_text(max_size=18, font_family='Calibri')
    print(f"   fit_text() called successfully")
    print(f"   Font size AFTER fit_text: {run2.font.size.pt if run2.font.size else 'None'}pt")
except Exception as e:
    print(f"   ERROR calling fit_text(): {e}")

# Test 3: Multiple runs (like our markdown)
print("\n3. Creating text box with MULTIPLE runs (markdown-style)...")
box3 = slide.shapes.add_textbox(Inches(7), Inches(1), Inches(2), Inches(1))
tf3 = box3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]

# Add multiple runs with different formatting (like markdown does)
run3a = p3.add_run()
run3a.text = "Bold text "
run3a.font.name = 'Calibri'
run3a.font.size = Pt(18)
run3a.font.bold = True

run3b = p3.add_run()
run3b.text = "and italic text "
run3b.font.name = 'Calibri'
run3b.font.size = Pt(18)
run3b.font.italic = True

run3c = p3.add_run()
run3c.text = "and more content to cause overflow"
run3c.font.name = 'Calibri'
run3c.font.size = Pt(18)

print(f"   Created 3 runs (bold, italic, normal)")
print(f"   Total text length: {len(tf3.text)} chars")

# Try fit_text on multiple runs
try:
    tf3.fit_text(max_size=18, font_family='Calibri')
    print(f"   fit_text() called successfully on multiple runs")
    print(f"   Run 1 size after: {run3a.font.size.pt if run3a.font.size else 'None'}pt")
    print(f"   Run 2 size after: {run3b.font.size.pt if run3b.font.size else 'None'}pt")
    print(f"   Run 3 size after: {run3c.font.size.pt if run3c.font.size else 'None'}pt")
except Exception as e:
    print(f"   ERROR calling fit_text(): {e}")

# Save
output_path = os.path.join(os.path.dirname(__file__), "..", "test_fittext_diagnostic.pptx")
prs.save(output_path)

print(f"\n" + "="*60)
print(f"Saved: {output_path}")
print(f"\nOPEN THE FILE AND CHECK:")
print(f"  Box 1 (left): Text should OVERFLOW (no fit_text)")
print(f"  Box 2 (middle): Text should FIT (fit_text applied)")
print(f"  Box 3 (right): Text should FIT (fit_text on multiple runs)")
print(f"\nIf boxes 2 and 3 still overflow, fit_text() is not working!")
print("="*60)
