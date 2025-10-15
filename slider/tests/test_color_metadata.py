"""
Test color metadata detection and application.
"""

import asyncio
import sys
import os
import base64

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_maker import slide_maker
from autohive_integrations_sdk import ExecutionContext
from pptx import Presentation
from io import BytesIO


async def test_color():
    auth = {}

    async with ExecutionContext(auth=auth) as context:
        result = await slide_maker.execute_action(
            "create_presentation",
            {"title": "Color Test"},
            context
        )

        presentation_id = result["presentation_id"]
        files = [result["file"]]

        # Add placeholder with color metadata
        await slide_maker.execute_action(
            "add_elements",
            {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "elements": [
                    {"content": "[Title, Color=#FF0000, Fontsize=32]", "position": {"left": 1, "top": 1, "width": 8, "height": 1}},
                    {"content": "[Subtitle, Color=#666666, Italic]", "position": {"left": 1, "top": 2.2, "width": 8, "height": 0.7}}
                ],
                "files": files
            },
            context
        )

        # Replace
        result = await slide_maker.execute_action(
            "find_and_replace",
            {
                "presentation_id": presentation_id,
                "replacements": [
                    {"find": "[Title, Color=#FF0000, Fontsize=32]", "replace": "Red Title"},
                    {"find": "[Subtitle, Color=#666666, Italic]", "replace": "Gray Subtitle"}
                ],
                "files": files
            },
            context
        )

        # Verify
        prs_bytes = base64.b64decode(result["file"]["content"])
        prs = Presentation(BytesIO(prs_bytes))
        slide = prs.slides[0]

        print('Color Application Test:')
        print()

        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]

                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                        rgb = run.font.color.rgb
                        hex_color = f'#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'
                        print(f'Element {idx}: \"{text}\"')
                        print(f'  Color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]}) = {hex_color}')
                        print(f'  Italic: {run.font.italic}')
                    else:
                        print(f'Element {idx}: \"{text}\" - No color')

        print()
        print('Expected:')
        print('  Element 0: Color=#FF0000 (Red)')
        print('  Element 1: Color=#666666 (Gray) + Italic=True')

asyncio.run(test_color())
