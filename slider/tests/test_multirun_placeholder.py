"""
Test multi-run placeholder issue (like [$x,xxx] split across runs in real template).
Creates placeholders that span multiple runs to simulate PowerPoint's behavior.
"""

import asyncio
import sys
import os
import base64

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_maker import slide_maker
from autohive_integrations_sdk import ExecutionContext
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO


async def test_multirun():
    """Test replacement when placeholder spans multiple runs"""

    print("=" * 70)
    print("MULTI-RUN PLACEHOLDER TEST")
    print("=" * 70)

    # Create a template with multi-run placeholder (simulate real template behavior)
    print("\n1. Creating template with multi-run placeholder...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add text box with placeholder split across runs
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]

    # Simulate how PowerPoint might split [$x,xxx] across runs
    # (This happens when users type/format differently in PowerPoint)
    run1 = p.add_run()
    run1.text = "Price: [$"
    run1.font.size = Pt(18)

    run2 = p.add_run()
    run2.text = "x,xxx"
    run2.font.size = Pt(18)
    run2.font.bold = True  # Different formatting causes split

    run3 = p.add_run()
    run3.text = "] per month"
    run3.font.size = Pt(18)

    # Save to bytes
    template_bytes = BytesIO()
    prs.save(template_bytes)
    template_bytes.seek(0)

    print("   [OK] Created template")
    print("   Placeholder structure:")
    print("     Run 0: 'Price: [$'")
    print("     Run 1: 'x,xxx' (bold)")
    print("     Run 2: '] per month'")
    print("   Full text: 'Price: [$x,xxx] per month'")

    # Test with integration
    auth = {}

    async with ExecutionContext(auth=auth) as context:
        template_b64 = base64.b64encode(template_bytes.getvalue()).decode('utf-8')
        files = [{
            "name": "multirun_template.pptx",
            "contentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "content": template_b64
        }]

        result = await slide_maker.execute_action(
            "create_presentation",
            {"files": files},
            context
        )

        presentation_id = result["presentation_id"]
        files = [result["file"]]

        print("\n2. Attempting to replace [$x,xxx]...")

        result = await slide_maker.execute_action(
            "find_and_replace",
            {
                "presentation_id": presentation_id,
                "replacements": [
                    {"find": "[$x,xxx]", "replace": "$5,000"}
                ],
                "files": files
            },
            context
        )

        print(f"\n   Status: {result['status']}")
        print(f"   Summary: {result['summary']}")

        if result['changes']:
            change = result['changes'][0]
            print(f"\n   Replacement status: {change['status']}")
            if change['status'] == 'not_found':
                print(f"   [PROBLEM] Placeholder spans runs - current code can't find it!")
            else:
                print(f"   [OK] Replacement worked!")

        # Verify actual text
        prs_result = Presentation(BytesIO(base64.b64decode(result["file"]["content"])))
        slide_result = prs_result.slides[0]
        final_text = slide_result.shapes[0].text_frame.text

        print(f"\n   Final text: \"{final_text}\"")

        if '[$x,xxx]' in final_text:
            print("   [PROBLEM] Placeholder still present - NOT replaced")
        elif '$5,000' in final_text:
            print("   [SUCCESS] Placeholder replaced correctly")
        else:
            print("   [UNKNOWN] Text changed but unclear if correct")

        print("\n" + "=" * 70)
        print("This confirms the multi-run placeholder issue!")
        print("=" * 70)


if __name__ == "__main__":
    asyncio.run(test_multirun())
