"""
Comprehensive test for placeholder detection and metadata application.
Tests both detection (get_slide_elements) and application (find_and_replace).
"""

import asyncio
import sys
import os
import base64
import json

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_maker import slide_maker
from autohive_integrations_sdk import ExecutionContext


async def test_placeholder_metadata():
    """Test placeholder detection and metadata application"""

    print("=" * 70)
    print("PLACEHOLDER METADATA DETECTION & APPLICATION TEST")
    print("=" * 70)

    auth = {}

    async with ExecutionContext(auth=auth) as context:
        # Create presentation
        result = await slide_maker.execute_action(
            "create_presentation",
            {"title": "Metadata Test"},
            context
        )

        presentation_id = result["presentation_id"]
        files = [result["file"]]

        print("\n1. Creating placeholders with various metadata...")

        # Add placeholders with different metadata
        await slide_maker.execute_action(
            "add_elements",
            {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "elements": [
                    {
                        "content": "[Title, Fontsize=32pt, Bold=true, Font=Sofia Pro]",
                        "position": {"left": 1, "top": 0.5, "width": 8, "height": 1}
                    },
                    {
                        "content": "[Subtitle, Fontsize=18pt, Italic=true, Color=#666666]",
                        "position": {"left": 1, "top": 1.7, "width": 8, "height": 0.7}
                    },
                    {
                        "content": "{Company Name}",  # Curly brace placeholder
                        "position": {"left": 1, "top": 2.6, "width": 4, "height": 0.6}
                    },
                    {
                        "content": "{{Date}}",  # Double curly placeholder
                        "position": {"left": 5.5, "top": 2.6, "width": 3.5, "height": 0.6}
                    },
                    {
                        "content": "[Body Text]",  # Simple placeholder (no metadata)
                        "position": {"left": 1, "top": 3.4, "width": 8, "height": 1.5}
                    },
                    {
                        "content": "Plain text with no placeholders",
                        "position": {"left": 1, "top": 5.1, "width": 8, "height": 0.6}
                    }
                ],
                "files": files
            },
            context
        )

        print("   [OK] Added 6 text boxes (5 with placeholders, 1 plain)")

        # Step 2: Call get_slide_elements to verify detection
        print("\n2. Testing placeholder detection (get_slide_elements)...")

        result = await slide_maker.execute_action(
            "get_slide_elements",
            {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "include_content": True,
                "files": files
            },
            context
        )

        print(f"\n   Total elements: {result['total_elements']}")
        print(f"\n   Placeholder Detection Results:")

        fillable_count = 0
        for elem in result['elements']:
            is_fillable = elem.get('is_fillable', False)
            if is_fillable:
                fillable_count += 1
                print(f"\n   Element {elem['index']}:")
                print(f"     Fillable: YES")
                print(f"     Placeholders: {elem.get('placeholders', [])}")
                if elem.get('placeholder_metadata'):
                    print(f"     Metadata: {json.dumps(elem['placeholder_metadata'], indent=8)}")
            else:
                print(f"\n   Element {elem['index']}: Not fillable")

        print(f"\n   [OK] Detected {fillable_count} fillable elements")

        # Step 3: Replace placeholders with metadata application
        print("\n3. Testing metadata application (find_and_replace)...")

        replacements = [
            {
                "find": "[Title, Fontsize=32pt, Bold=true, Font=Sofia Pro]",
                "replace": "Digital Transformation Strategy"
                # Metadata will be auto-extracted and applied!
            },
            {
                "find": "[Subtitle, Fontsize=18pt, Italic=true, Color=#666666]",
                "replace": "A comprehensive roadmap for modernization"
            },
            {
                "find": "{Company Name}",
                "replace": "TechVenture Solutions"
            },
            {
                "find": "{{Date}}",
                "replace": "September 30, 2025"
            },
            {
                "find": "[Body Text]",
                "replace": "This is body text that will be auto-sized based on available space."
            }
        ]

        result = await slide_maker.execute_action(
            "find_and_replace",
            {
                "presentation_id": presentation_id,
                "replacements": replacements,
                "files": files
            },
            context
        )

        print(f"\n   Status: {result['status']}")
        print(f"   Summary: {result['summary']['successful']}/{result['summary']['requested']} successful")

        if result.get('changes'):
            print(f"\n   Metadata Application Results:")
            for change in result['changes']:
                if change['status'] == 'replaced':
                    print(f"\n   {change['find'][:50]}...")
                    print(f"     Status: {change['status']}")
                    print(f"     Font size: {change.get('font_size_applied')}pt")
                    if change.get('forced'):
                        print(f"     Source: Forced parameter")
                    elif 'Fontsize=' in change['find']:
                        print(f"     Source: Placeholder metadata")

        files = [result["file"]]

        # Step 4: Verify formatting was applied
        print("\n4. Verifying applied formatting...")

        from pptx import Presentation
        from io import BytesIO

        prs_bytes = base64.b64decode(files[0]["content"])
        prs = Presentation(BytesIO(prs_bytes))
        slide = prs.slides[0]

        print(f"\n   Checking applied formatting:")

        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text[:35]
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    print(f"\n   Element {idx}: \"{text}...\"")
                    print(f"     Size: {run.font.size.pt if run.font.size else 'None'}pt")
                    print(f"     Bold: {run.font.bold}")
                    print(f"     Italic: {run.font.italic}")
                    print(f"     Font: {run.font.name}")
                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                        rgb = run.font.color.rgb
                        print(f"     Color: rgb({rgb[0]}, {rgb[1]}, {rgb[2]})")

        # Save output
        import time
        output_path = os.path.join(
            os.path.dirname(__file__),
            f"test_metadata_{int(time.time())}.pptx"
        )

        with open(output_path, 'wb') as f:
            f.write(prs_bytes)

        print(f"\n5. Saved to: {output_path}")

        print("\n" + "=" * 70)
        print("[SUCCESS] PLACEHOLDER METADATA TEST COMPLETE!")
        print("=" * 70)
        print("\nEXPECTED RESULTS:")
        print("  Element 0: 32pt, Bold, Sofia Pro font (from metadata)")
        print("  Element 1: 18pt, Italic, Gray color (from metadata)")
        print("  Element 2-4: Auto-sized (no metadata)")
        print(f"\n[FILE] Output: {output_path}")
        print("=" * 70)


if __name__ == "__main__":
    asyncio.run(test_placeholder_metadata())
