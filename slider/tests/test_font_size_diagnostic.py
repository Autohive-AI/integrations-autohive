"""
Diagnostic tool to analyze font sizes in a presentation.
Shows actual font size vs requested size for all text elements.
"""

import sys
import os
import base64

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from io import BytesIO


def analyze_font_sizes(pptx_path):
    """Analyze all font sizes in a presentation"""

    print("=" * 70)
    print("FONT SIZE DIAGNOSTIC TOOL")
    print("=" * 70)

    # Load presentation
    prs = Presentation(pptx_path)

    print(f"\nPresentation: {os.path.basename(pptx_path)}")
    print(f"Total slides: {len(prs.slides)}")

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n{'='*70}")
        print(f"SLIDE {slide_idx}")
        print(f"{'='*70}")

        if len(slide.shapes) == 0:
            print("  (No elements)")
            continue

        for shape_idx, shape in enumerate(slide.shapes):
            # Check if shape has text
            if not hasattr(shape, "text_frame"):
                continue

            text_frame = shape.text_frame
            if not text_frame.text or len(text_frame.text.strip()) == 0:
                continue

            # Get shape dimensions
            width_in = shape.width / 914400
            height_in = shape.height / 914400
            left_in = shape.left / 914400
            top_in = shape.top / 914400

            # Get text content (truncated)
            text_preview = text_frame.text[:50].replace('\n', ' ')
            if len(text_frame.text) > 50:
                text_preview += "..."

            print(f"\nElement {shape_idx}: {shape.shape_type}")
            print(f"  Position: ({left_in:.2f}, {top_in:.2f})")
            print(f"  Size: {width_in:.2f}w x {height_in:.2f}h inches")
            print(f"  Text: \"{text_preview}\"")
            print(f"  Text length: {len(text_frame.text)} chars")

            # Analyze font sizes in paragraphs
            font_sizes = []
            for para in text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)

            if font_sizes:
                min_size = min(font_sizes)
                max_size = max(font_sizes)
                avg_size = sum(font_sizes) / len(font_sizes)

                print(f"  Font sizes: min={min_size:.0f}pt, max={max_size:.0f}pt, avg={avg_size:.1f}pt")

                # Highlight potential issues
                if max_size > 20:
                    print(f"  [!] POTENTIALLY TOO BIG: {max_size:.0f}pt")
                if min_size < 10:
                    print(f"  [!] POTENTIALLY TOO SMALL: {min_size:.0f}pt")

                # Check text density
                chars_per_sq_inch = len(text_frame.text) / (width_in * height_in)
                if chars_per_sq_inch > 100:
                    print(f"  [!] HIGH DENSITY: {chars_per_sq_inch:.0f} chars/sq.in (may overflow)")
                elif chars_per_sq_inch < 20:
                    print(f"  [!] LOW DENSITY: {chars_per_sq_inch:.0f} chars/sq.in (could be bigger)")
            else:
                print(f"  Font sizes: No font size set")

    print("\n" + "=" * 70)


if __name__ == "__main__":
    import sys

    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        # Use latest test output
        test_dir = os.path.dirname(__file__)
        files = [f for f in os.listdir(test_dir) if f.startswith("test_agent_replication_") and f.endswith(".pptx")]
        if files:
            files.sort(reverse=True)  # Latest first
            pptx_path = os.path.join(test_dir, files[0])
        else:
            print("No test file found. Usage: python test_font_size_diagnostic.py <path_to_pptx>")
            sys.exit(1)

    if os.path.exists(pptx_path):
        analyze_font_sizes(pptx_path)
    else:
        print(f"File not found: {pptx_path}")
