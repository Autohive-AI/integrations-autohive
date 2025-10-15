"""
Check what auto-sizing is configured on the HK Template boxes.
This will reveal if boxes have TEXT_TO_FIT_SHAPE enabled (causing resizing).
"""

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
import os

template_path = os.path.join(
    os.path.dirname(__file__),
    "..",
    "HK Template.pptx"
)

print("=" * 70)
print("TEMPLATE AUTO-SIZE DIAGNOSTIC")
print("=" * 70)

prs = Presentation(template_path)

print(f"\nTemplate: {os.path.basename(template_path)}")
print(f"Total slides: {len(prs.slides)}")

autosize_map = {
    MSO_AUTO_SIZE.NONE: "NONE",
    MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT: "SHAPE_TO_FIT_TEXT (RESIZES BOX!)",
    MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE: "TEXT_TO_FIT_SHAPE (SHRINKS FONT!)"
}

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n{'='*70}")
    print(f"SLIDE {slide_idx}")
    print(f"{'='*70}")

    for shape_idx, shape in enumerate(slide.shapes):
        if not hasattr(shape, "text_frame"):
            continue

        text_frame = shape.text_frame
        if not text_frame.text or len(text_frame.text.strip()) == 0:
            continue

        text_preview = text_frame.text[:40].replace('\n', ' ')
        if len(text_frame.text) > 40:
            text_preview += "..."

        autosize = text_frame.auto_size
        autosize_str = autosize_map.get(autosize, f"UNKNOWN ({autosize})")

        print(f"\nElement {shape_idx}:")
        print(f"  Text: \"{text_preview}\"")
        print(f"  Auto-size: {autosize_str}")

        if autosize == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
            print(f"  [!] WARNING: Box will RESIZE to fit text!")
            print(f"      This is likely causing Slide 6 snapping behavior")
        elif autosize == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
            print(f"  [!] WARNING: Font will SHRINK to fit box!")
            print(f"      This will override our calculated font size")

print("\n" + "=" * 70)
print("KEY FINDINGS:")
print("=" * 70)
print("\nIf any boxes have SHAPE_TO_FIT_TEXT:")
print("  -> That's why Slide 6 box is resizing itself")
print("  -> find_and_replace should DISABLE this after replacing text")
print("\nIf any boxes have TEXT_TO_FIT_SHAPE:")
print("  -> PowerPoint will override our calculated font size")
print("  -> This could cause overflow if PP calculates differently")
print("=" * 70)
