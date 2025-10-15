"""
Test placeholder metadata with real template workflow.
1. Load template with metadata placeholders
2. Detect placeholders with get_slide_elements
3. Fill with find_and_replace
4. Verify all metadata formatting applied
"""

import asyncio
import sys
import os
import base64
import json

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_maker import slide_maker
from autohive_integrations_sdk import ExecutionContext
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO


async def test_real_workflow():
    """Test with manually created template"""

    print("=" * 70)
    print("REAL TEMPLATE WORKFLOW TEST")
    print("=" * 70)

    # Step 1: Create a template manually
    print("\n1. Creating template with metadata placeholders...")

    from pptx import Presentation as PptxPresentation
    prs = PptxPresentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add blank slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add placeholder text boxes
    # Placeholder 1: With full metadata
    box1 = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf1 = box1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "[Title, Fontsize=32pt, Bold=true]"
    p1.runs[0].font.size = Pt(14)  # Template default

    # Placeholder 2: With italic and color
    box2 = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(8), Inches(0.7))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "[Subtitle, Fontsize=18pt, Italic=true]"
    p2.runs[0].font.size = Pt(12)

    # Placeholder 3: Simple (no metadata)
    box3 = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(8), Inches(1))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "[Company]"
    p3.runs[0].font.size = Pt(16)

    # Save template to bytes
    template_bytes = BytesIO()
    prs.save(template_bytes)
    template_bytes.seek(0)
    template_b64 = base64.b64encode(template_bytes.getvalue()).decode('utf-8')

    print("   [OK] Template created with 3 placeholders")

    # Step 2: Test with integration
    auth = {}

    async with ExecutionContext(auth=auth) as context:
        files = [{
            "name": "template.pptx",
            "contentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "content": template_b64
        }]

        result = await slide_maker.execute_action(
            "create_presentation",
            {"files": files},
            context
        )

        presentation_id = result["presentation_id"]
        files = [result["file"]]

        # Step 3: Detect placeholders
        print("\n2. Detecting placeholders...")

        result = await slide_maker.execute_action(
            "get_slide_elements",
            {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "include_content": True,
                "files": files
            },
            context
        )

        print(f"\n   Detection results:")
        for elem in result['elements']:
            if elem.get('is_fillable'):
                print(f"\n   Element {elem['index']}:")
                print(f"     Content: {elem.get('content', '')[:50]}")
                print(f"     Placeholders: {elem.get('placeholders')}")
                if elem.get('placeholder_metadata'):
                    print(f"     Metadata: {elem.get('placeholder_metadata')}")

        # Step 4: Fill placeholders
        print("\n3. Filling placeholders with metadata application...")

        replacements = [
            {"find": "[Title, Fontsize=32pt, Bold=true]", "replace": "Q4 Business Review"},
            {"find": "[Subtitle, Fontsize=18pt, Italic=true]", "replace": "Executive Summary"},
            {"find": "[Company]", "replace": "TechVenture Solutions Ltd"}
        ]

        result = await slide_maker.execute_action(
            "find_and_replace",
            {
                "presentation_id": presentation_id,
                "replacements": replacements,
                "files": files
            },
            context
        )

        print(f"\n   Status: {result['status']}")
        print(f"   Replacements: {result['summary']['successful']}/{result['summary']['requested']}")

        files = [result["file"]]

        # Step 5: Verify formatting
        print("\n4. Verifying applied formatting...")

        prs_bytes = base64.b64decode(files[0]["content"])
        prs_result = Presentation(BytesIO(prs_bytes))
        slide_result = prs_result.slides[0]

        print(f"\n   Formatting verification:")
        for idx, shape in enumerate(slide_result.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text[:30]
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    print(f"\n   Element {idx}: \"{text}...\"")
                    print(f"     Size: {run.font.size.pt if run.font.size else 'N/A'}pt")
                    print(f"     Bold: {run.font.bold}")
                    print(f"     Italic: {run.font.italic}")

        # Save
        output_path = os.path.join(os.path.dirname(__file__), f"test_real_workflow_{int(__import__('time').time())}.pptx")
        with open(output_path, 'wb') as f:
            f.write(prs_bytes)

        print(f"\n5. Saved to: {output_path}")

        print("\n" + "=" * 70)
        print("EXPECTED:")
        print("  Element 0: 32pt, Bold=True (from [Title, Fontsize=32pt, Bold=true])")
        print("  Element 1: 18pt, Italic=True (from [Subtitle, Fontsize=18pt, Italic=true])")
        print("  Element 2: Auto-sized (no metadata)")
        print("=" * 70)


if __name__ == "__main__":
    asyncio.run(test_real_workflow())
