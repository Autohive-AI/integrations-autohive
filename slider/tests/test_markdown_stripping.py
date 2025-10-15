"""
Test that markdown is stripped when placeholder metadata specifies formatting.
Prevents literal **/** characters appearing when formatting is applied via metadata.
"""

import asyncio
import sys
import os
import base64

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_maker import slide_maker
from autohive_integrations_sdk import ExecutionContext
from pptx import Presentation
from io import BytesIO


async def test_markdown_stripping():
    """Test markdown stripping with metadata"""

    print("=" * 70)
    print("MARKDOWN STRIPPING TEST")
    print("=" * 70)

    auth = {}

    async with ExecutionContext(auth=auth) as context:
        # Create presentation with metadata placeholders
        result = await slide_maker.execute_action(
            "create_presentation",
            {"title": "Markdown Stripping Test"},
            context
        )

        presentation_id = result["presentation_id"]
        files = [result["file"]]

        # Add placeholders WITH metadata
        await slide_maker.execute_action(
            "add_elements",
            {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "elements": [
                    {
                        "content": "[Title, Bold=true, Fontsize=24pt]",
                        "position": {"left": 1, "top": 1, "width": 8, "height": 1}
                    },
                    {
                        "content": "[Subtitle, Italic=true, Fontsize=18pt]",
                        "position": {"left": 1, "top": 2.2, "width": 8, "height": 0.8}
                    },
                    {
                        "content": "[Body]",  # No metadata
                        "position": {"left": 1, "top": 3.2, "width": 8, "height": 1}
                    }
                ],
                "files": files
            },
            context
        )

        print("\n1. Created placeholders with metadata")

        # Replace with markdown in the replacement text (should be stripped)
        print("\n2. Replacing with markdown formatting in replacement text...")
        print("   (Markdown should be stripped where metadata exists)")

        replacements = [
            {
                "find": "[Title, Bold=true, Fontsize=24pt]",
                "replace": "**Q4 Business Review**"  # ← Should strip ** (Bold=true in metadata)
            },
            {
                "find": "[Subtitle, Italic=true, Fontsize=18pt]",
                "replace": "*Executive Summary*"  # ← Should strip * (Italic=true in metadata)
            },
            {
                "find": "[Body]",
                "replace": "This has **bold** and *italic* text"  # ← Should KEEP (no metadata)
            }
        ]

        result = await slide_maker.execute_action(
            "find_and_replace",
            {
                "presentation_id": presentation_id,
                "replacements": replacements,
                "files": files
            },
            context
        )

        print(f"\n   Status: {result['status']}")

        files = [result["file"]]

        # Verify stripping
        print("\n3. Verifying markdown was stripped correctly...")

        prs_bytes = base64.b64decode(files[0]["content"])
        prs = Presentation(BytesIO(prs_bytes))
        slide = prs.slides[0]

        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    print(f"\n   Element {idx}:")
                    print(f"     Text: \"{text}\"")
                    print(f"     Bold: {run.font.bold}")
                    print(f"     Italic: {run.font.italic}")

                    # Check for literal markdown characters
                    if '**' in text or '__' in text:
                        print(f"     [WARN] Contains literal ** or __ characters!")
                    if text.count('*') > 0 and text.count('*') % 2 == 0:
                        print(f"     [WARN] May contain literal * characters!")

        print("\n" + "=" * 70)
        print("EXPECTED RESULTS:")
        print("=" * 70)
        print('\n  Element 0: "Q4 Business Review" (NO ** visible)')
        print("            Bold: True (from metadata)")
        print('\n  Element 1: "Executive Summary" (NO * visible)')
        print("            Italic: True (from metadata)")
        print('\n  Element 2: "This has **bold** and *italic* text" (** and * VISIBLE)')
        print("            Bold/Italic: None (no metadata, markdown preserved)")
        print("=" * 70)


if __name__ == "__main__":
    asyncio.run(test_markdown_stripping())
