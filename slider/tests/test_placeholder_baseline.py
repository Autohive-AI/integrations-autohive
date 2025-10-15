"""
Baseline test: Verify which placeholder patterns currently work.
This ensures we don't break working functionality when fixing multi-run issue.
"""

import asyncio
import sys
import os
import base64

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_maker import slide_maker
from autohive_integrations_sdk import ExecutionContext


async def test_baseline():
    """Test various placeholder patterns to see what currently works"""

    print("=" * 70)
    print("PLACEHOLDER BASELINE TEST")
    print("=" * 70)
    print("\nTesting which placeholder patterns currently work for replacement...")

    auth = {}

    async with ExecutionContext(auth=auth) as context:
        result = await slide_maker.execute_action(
            "create_presentation",
            {"title": "Baseline Test"},
            context
        )

        presentation_id = result["presentation_id"]
        files = [result["file"]]

        # Create various placeholder patterns
        test_placeholders = [
            {"content": "[SIMPLE]", "label": "Simple"},
            {"content": "[TWO WORDS]", "label": "Two words"},
            {"content": "[Title, Bold]", "label": "With metadata"},
            {"content": "{CURLY}", "label": "Curly braces"},
            {"content": "[$x,xxx]", "label": "Dollar and commas"},
            {"content": "[email@example.com]", "label": "Email format"},
            {"content": "[Date / 13 October 2025]", "label": "Date with slashes"},
        ]

        print("\n1. Creating placeholders...")

        elements = []
        for i, test in enumerate(test_placeholders):
            elements.append({
                "content": test["content"],
                "position": {"left": 1, "top": 1 + (i * 0.7), "width": 6, "height": 0.6}
            })

        await slide_maker.execute_action(
            "add_elements",
            {
                "presentation_id": presentation_id,
                "slide_index": 0,
                "elements": elements,
                "files": files
            },
            context
        )

        print(f"   [OK] Created {len(test_placeholders)} test placeholders")

        # Try to replace each one
        print("\n2. Testing replacements...")

        replacements = [
            {"find": "[SIMPLE]", "replace": "Replaced Simple"},
            {"find": "[TWO WORDS]", "replace": "Replaced Two Words"},
            {"find": "[Title, Bold]", "replace": "Replaced Metadata"},
            {"find": "{CURLY}", "replace": "Replaced Curly"},
            {"find": "[$x,xxx]", "replace": "$5,000"},
            {"find": "[email@example.com]", "replace": "test@test.com"},
            {"find": "[Date / 13 October 2025]", "replace": "30 Sept 2025"},
        ]

        result = await slide_maker.execute_action(
            "find_and_replace",
            {
                "presentation_id": presentation_id,
                "replacements": replacements,
                "files": files
            },
            context
        )

        print(f"\n   Status: {result['status']}")
        print(f"   Summary: {result['summary']}")

        # Show which worked and which didn't
        print("\n3. Results per placeholder:")
        for change in result.get('changes', []):
            status_icon = "OK" if change['status'] == 'replaced' else "FAIL"
            print(f"   [{status_icon}] {change['find']}: {change['status']}")
            if change['status'] == 'not_found':
                print(f"       Suggestion: {change.get('suggestion', 'N/A')}")

        # Verify in actual file
        from pptx import Presentation
        from io import BytesIO

        prs_bytes = base64.b64decode(result["file"]["content"])
        prs = Presentation(BytesIO(prs_bytes))
        slide = prs.slides[0]

        print("\n4. Verification (actual text in output):")
        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text
                print(f"   Element {idx}: \"{text}\"")

        print("\n" + "=" * 70)
        print("BASELINE ESTABLISHED")
        print("=" * 70)
        print("\nUse this to verify no working placeholders break after fix!")
        print("=" * 70)


if __name__ == "__main__":
    asyncio.run(test_baseline())
