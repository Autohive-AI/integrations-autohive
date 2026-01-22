import pytest
from unittest.mock import AsyncMock, MagicMock, patch
import io


@pytest.fixture
def mock_context():
    context = MagicMock()
    context.fetch = AsyncMock()
    return context


class TestListPresentationsAction:
    @pytest.mark.asyncio
    async def test_list_presentations_success(self, mock_context):
        from microsoft_powerpoint import ListPresentationsAction

        mock_context.fetch.return_value = {
            'value': [
                {
                    'id': 'item-123',
                    'name': 'Quarterly Report.pptx',
                    'webUrl': 'https://example.com/file',
                    'lastModifiedDateTime': '2024-01-15T10:30:00Z',
                    'size': 2048576
                }
            ],
            '@odata.nextLink': None
        }

        action = ListPresentationsAction()
        result = await action.execute({}, mock_context)

        assert result['result'] is True
        assert len(result['presentations']) == 1
        assert result['presentations'][0]['id'] == 'item-123'
        assert result['presentations'][0]['name'] == 'Quarterly Report.pptx'

    @pytest.mark.asyncio
    async def test_list_presentations_with_filter(self, mock_context):
        from microsoft_powerpoint import ListPresentationsAction

        mock_context.fetch.return_value = {'value': [], '@odata.nextLink': None}

        action = ListPresentationsAction()
        result = await action.execute({
            'name_contains': 'Report',
            'folder_path': 'Documents',
            'page_size': 10
        }, mock_context)

        assert result['result'] is True
        mock_context.fetch.assert_called_once()

    @pytest.mark.asyncio
    async def test_list_presentations_escapes_quotes(self, mock_context):
        from microsoft_powerpoint import ListPresentationsAction

        mock_context.fetch.return_value = {'value': [], '@odata.nextLink': None}

        action = ListPresentationsAction()
        result = await action.execute({
            'name_contains': "O'Brien's Report"
        }, mock_context)

        assert result['result'] is True
        call_args = mock_context.fetch.call_args
        params = call_args.kwargs.get('params', {})
        assert "O''Brien''s Report" in params.get('$filter', '')

    @pytest.mark.asyncio
    async def test_list_presentations_error(self, mock_context):
        from microsoft_powerpoint import ListPresentationsAction

        mock_context.fetch.side_effect = Exception("API Error")

        action = ListPresentationsAction()
        result = await action.execute({}, mock_context)

        assert result['result'] is False
        assert 'error' in result
        assert 'API Error' in result['error']


class TestGetPresentationAction:
    @pytest.mark.asyncio
    async def test_get_presentation_success(self, mock_context):
        from microsoft_powerpoint import GetPresentationAction

        mock_context.fetch.return_value = {
            'id': 'item-123',
            'name': 'Quarterly Report.pptx',
            'size': 2048576,
            'webUrl': 'https://example.com/file',
            'createdDateTime': '2024-01-10T08:00:00Z',
            'lastModifiedDateTime': '2024-01-15T10:30:00Z',
            'createdBy': {'user': {'displayName': 'John Doe'}},
            'lastModifiedBy': {'user': {'displayName': 'Jane Doe'}},
            'parentReference': {'path': '/drive/root:/Documents'}
        }

        action = GetPresentationAction()
        result = await action.execute({'presentation_id': 'item-123'}, mock_context)

        assert result['result'] is True
        assert result['id'] == 'item-123'
        assert result['name'] == 'Quarterly Report.pptx'
        assert result['size'] == 2048576


class TestGetSlidesAction:
    @pytest.mark.asyncio
    @patch('microsoft_powerpoint.download_file_content')
    async def test_get_slides_success(self, mock_download, mock_context):
        from microsoft_powerpoint import GetSlidesAction
        from pptx import Presentation
        
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.slides.add_slide(prs.slide_layouts[6])
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_download.return_value = buffer.read()

        mock_context.fetch.return_value = {
            'value': [{'id': '0', 'medium': {'url': 'https://thumb.com', 'width': 800, 'height': 600}}]
        }

        action = GetSlidesAction()
        result = await action.execute({'presentation_id': 'item-123'}, mock_context)

        assert result['result'] is True
        assert result['slide_count'] == 2
        assert len(result['slides']) == 2
        assert result['slides'][0]['index'] == 1


class TestGetSlideAction:
    @pytest.mark.asyncio
    @patch('microsoft_powerpoint.download_file_content')
    async def test_get_slide_success(self, mock_download, mock_context):
        from microsoft_powerpoint import GetSlideAction
        from pptx import Presentation
        
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.slides.add_slide(prs.slide_layouts[6])
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_download.return_value = buffer.read()

        mock_context.fetch.return_value = {
            'value': [{'id': '0', 'large': {'url': 'https://thumb.com', 'width': 1600, 'height': 1200}}]
        }

        action = GetSlideAction()
        result = await action.execute({
            'presentation_id': 'item-123',
            'slide_index': 2
        }, mock_context)

        assert result['result'] is True
        assert result['index'] == 2

    @pytest.mark.asyncio
    async def test_get_slide_invalid_index(self, mock_context):
        from microsoft_powerpoint import GetSlideAction

        action = GetSlideAction()
        result = await action.execute({
            'presentation_id': 'item-123',
            'slide_index': 0
        }, mock_context)

        assert result['result'] is False
        assert 'error' in result

    @pytest.mark.asyncio
    @patch('microsoft_powerpoint.download_file_content')
    async def test_get_slide_out_of_range(self, mock_download, mock_context):
        from microsoft_powerpoint import GetSlideAction
        from pptx import Presentation
        
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_download.return_value = buffer.read()

        action = GetSlideAction()
        result = await action.execute({
            'presentation_id': 'item-123',
            'slide_index': 5
        }, mock_context)

        assert result['result'] is False
        assert 'out of range' in result['error']


class TestDeleteSlideAction:
    @pytest.mark.asyncio
    @patch('microsoft_powerpoint.overwrite_file_content')
    @patch('microsoft_powerpoint.download_file_content')
    async def test_delete_slide_success(self, mock_download, mock_overwrite, mock_context):
        from microsoft_powerpoint import DeleteSlideAction
        from pptx import Presentation
        
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.slides.add_slide(prs.slide_layouts[6])
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_download.return_value = buffer.read()
        mock_overwrite.return_value = {}

        action = DeleteSlideAction()
        result = await action.execute({
            'presentation_id': 'item-123',
            'slide_index': 1
        }, mock_context)

        assert result['result'] is True
        assert result['deleted'] is True
        assert result['slide_count'] == 1

    @pytest.mark.asyncio
    @patch('microsoft_powerpoint.download_file_content')
    async def test_delete_last_slide_error(self, mock_download, mock_context):
        from microsoft_powerpoint import DeleteSlideAction
        from pptx import Presentation
        
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_download.return_value = buffer.read()

        action = DeleteSlideAction()
        result = await action.execute({
            'presentation_id': 'item-123',
            'slide_index': 1
        }, mock_context)

        assert result['result'] is False
        assert 'Cannot delete the last slide' in result['error']


class TestExportPdfAction:
    @pytest.mark.asyncio
    async def test_export_pdf_success(self, mock_context):
        from microsoft_powerpoint import ExportPdfAction

        mock_context.fetch.side_effect = [
            {'name': 'Report.pptx', 'parentReference': {'path': '/drive/root:/Documents'}},
            b'%PDF-1.4...',
            {'id': 'pdf-123', 'name': 'Report.pdf', 'webUrl': 'https://example.com/pdf', 'size': 1024},
            {'@microsoft.graph.downloadUrl': 'https://download.com/pdf'}
        ]

        action = ExportPdfAction()
        result = await action.execute({'presentation_id': 'item-123'}, mock_context)

        assert result['result'] is True
        assert result['pdf_id'] == 'pdf-123'
        assert result['pdf_name'] == 'Report.pdf'


class TestGetSlideImageAction:
    @pytest.mark.asyncio
    @patch('microsoft_powerpoint.download_file_content')
    async def test_get_slide_image_success(self, mock_download, mock_context):
        from microsoft_powerpoint import GetSlideImageAction
        from pptx import Presentation
        
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_download.return_value = buffer.read()

        mock_context.fetch.return_value = {
            'value': [{'id': '0', 'large': {'url': 'https://image.com/slide1.png', 'width': 1600, 'height': 1200}}]
        }

        action = GetSlideImageAction()
        result = await action.execute({
            'presentation_id': 'item-123',
            'slide_index': 1
        }, mock_context)

        assert result['result'] is True
        assert result['image_url'] == 'https://image.com/slide1.png'
        assert result['format'] == 'png'
        assert 'note' in result

    @pytest.mark.asyncio
    async def test_get_slide_image_invalid_index(self, mock_context):
        from microsoft_powerpoint import GetSlideImageAction

        action = GetSlideImageAction()
        result = await action.execute({
            'presentation_id': 'item-123',
            'slide_index': 0
        }, mock_context)

        assert result['result'] is False


class TestCreatePresentationAction:
    @pytest.mark.asyncio
    async def test_create_presentation_from_template(self, mock_context):
        from microsoft_powerpoint import CreatePresentationAction

        mock_context.fetch.return_value = {
            'id': 'new-item-123',
            'name': 'New Presentation.pptx',
            'webUrl': 'https://example.com/new',
            'createdDateTime': '2024-01-20T10:00:00Z'
        }

        action = CreatePresentationAction()
        result = await action.execute({
            'name': 'New Presentation',
            'template_id': 'template-123'
        }, mock_context)

        assert result['result'] is True

    @pytest.mark.asyncio
    @patch('microsoft_powerpoint.create_blank_pptx')
    @patch('microsoft_powerpoint.upload_file_content')
    async def test_create_blank_presentation(self, mock_upload, mock_create_blank, mock_context):
        from microsoft_powerpoint import CreatePresentationAction

        mock_create_blank.return_value = b'PK...'
        mock_upload.return_value = {
            'id': 'new-item-123',
            'name': 'Blank.pptx',
            'webUrl': 'https://example.com/blank',
            'createdDateTime': '2024-01-20T10:00:00Z'
        }

        action = CreatePresentationAction()
        result = await action.execute({'name': 'Blank'}, mock_context)

        assert result['result'] is True
        mock_create_blank.assert_called_once()


class TestHelperFunctions:
    def test_odata_escape(self):
        from microsoft_powerpoint import odata_escape
        
        assert odata_escape("normal") == "normal"
        assert odata_escape("O'Brien") == "O''Brien"
        assert odata_escape("It's a test's") == "It''s a test''s"
